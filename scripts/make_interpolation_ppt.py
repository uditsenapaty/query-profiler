#!/usr/bin/env python3
# =========================================================
# scripts/make_interpolation_ppt.py
# ---------------------------------------------------------
# Builds  interpolation_methods_explained.pptx  — a formula-
# first deck explaining every sampling/interpolation method
# in scripts/interpolation.py, judged by every metric in
# scripts/evaluate_interpolation.py, on a worked 5x5 grid,
# closing with the real Q5 results for parallel ON (s1q2)
# and parallel OFF (s1q0).
#
# Numbers come from scripts/interp_ppt_data.py (exact math).
# Run:  python3 scripts/make_interpolation_ppt.py
# =========================================================

import os
import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

import interp_ppt_data as T

# ---------------------------------------------------------
# palette / fonts
# ---------------------------------------------------------
INK    = RGBColor(0x1F, 0x2A, 0x37)   # near-black slate
ACCENT = RGBColor(0x0E, 0x7C, 0x86)   # teal
ACC2   = RGBColor(0xB4, 0x53, 0x09)   # amber (winner / highlight)
GREEN  = RGBColor(0x15, 0x6F, 0x3B)
RED    = RGBColor(0xB1, 0x2A, 0x2A)
GREY   = RGBColor(0x6B, 0x72, 0x80)
LIGHT  = RGBColor(0xF2, 0xF4, 0xF6)   # formula box bg
ROW1   = RGBColor(0xFF, 0xFF, 0xFF)
ROW2   = RGBColor(0xEC, 0xF3, 0xF4)   # light teal band
HIROW  = RGBColor(0xFD, 0xF1, 0xDD)   # amber band (winner row)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BAND   = RGBColor(0xE8, 0xEC, 0xEF)

FT  = "Segoe UI"
FM  = "Consolas"

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]

_page = [0]

# ---------------------------------------------------------
# low-level helpers
# ---------------------------------------------------------
def _set_fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def _runs(p, runs, default_size, default_color, font=FT, align=None):
    """runs = str OR list of (text, dict-opts)."""
    if align is not None:
        p.alignment = align
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, o in runs:
        r = p.add_run(); r.text = text
        f = r.font
        f.name  = o.get("font", font)
        f.size  = Pt(o.get("size", default_size))
        f.bold  = o.get("bold", False)
        f.italic= o.get("italic", False)
        f.color.rgb = o.get("color", default_color)
    return p

def textbox(slide, l, t, w, h, lines, size=16, color=INK, font=FT,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=4,
            line_spacing=1.0):
    """lines = list; each item is str | (text,opts) | list-of-(text,opts)."""
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        _runs(p, ln, size, color, font=font, align=align)
    return tb

def rrect(slide, l, t, w, h, fill, line=None, line_w=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sp.adjustments[0] = 0.06
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def rect(slide, l, t, w, h, fill):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    _set_fill(sp, fill); sp.shadow.inherit = False
    return sp

# ---------------------------------------------------------
# slide scaffold
# ---------------------------------------------------------
def slide(kicker, title, accent=ACCENT):
    s = prs.slides.add_slide(BLANK)
    _page[0] += 1
    rect(s, 0, 0, SW, Inches(0.16), accent)                       # top hairline
    # kicker
    textbox(s, Inches(0.55), Inches(0.24), Inches(11.5), Inches(0.3),
            [[(kicker.upper(), {"size": 12.5, "bold": True, "color": accent})]],
            )
    # title
    textbox(s, Inches(0.55), Inches(0.52), Inches(12.3), Inches(0.8),
            [[(title, {"size": 27, "bold": True, "color": INK})]])
    rect(s, Inches(0.55), Inches(1.26), Inches(12.23), Pt(1.6), BAND)
    # footer
    textbox(s, Inches(0.55), Inches(7.12), Inches(9), Inches(0.3),
            [[("Interpolating Runtime Q-errors  ·  scripts/interpolation.py", {"size": 9, "color": GREY})]])
    textbox(s, Inches(12.2), Inches(7.12), Inches(0.9), Inches(0.3),
            [[(str(_page[0]), {"size": 9, "color": GREY})]], align=PP_ALIGN.RIGHT)
    return s

def mono(slide, l, t, w, h, lines, size=14, title=None, fill=LIGHT,
         color=INK, accent=ACCENT, line_spacing=1.05):
    """Formula / code panel: light rounded box + monospace text."""
    box = rrect(slide, l, t, w, h, fill, line=BAND, line_w=1.0)
    pad = Inches(0.14)
    ty = t + Inches(0.06)
    if title:
        textbox(slide, l + pad, ty, w - 2*pad, Inches(0.3),
                [[(title, {"size": size-1.5, "bold": True, "color": accent, "font": FT})]])
        ty = ty + Inches(0.34)
    textbox(slide, l + pad, ty, w - 2*pad, h - (ty - t) - Inches(0.06),
            lines, size=size, color=color, font=FM, line_spacing=line_spacing,
            space_after=3)
    return box

def bullets(slide, l, t, w, h, items, size=15.5, color=INK, accent=ACCENT,
            gap=5, line_spacing=1.02):
    """items = list of (level, runs) ; runs = str|list. level 0/1/2."""
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True
    marks = {0: "▸  ", 1: "–  ", 2: "·  "}
    for i, it in enumerate(items):
        lvl, runs = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(gap); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        mk = marks.get(lvl, "")
        mcol = accent if lvl == 0 else (GREY if lvl == 2 else ACCENT)
        runs2 = [(mk, {"color": mcol, "bold": True})]
        if isinstance(runs, str):
            runs2.append((runs, {}))
        else:
            runs2.extend(runs)
        _runs(p, runs2, size if lvl == 0 else size-1.5, color)
    return tb

def _cell(cell, runs, size, color, bold, align, fill, font=FT):
    cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    if isinstance(runs, str): runs = [(runs, {})]
    for text, o in runs:
        r = p.add_run(); r.text = text; f = r.font
        f.name = o.get("font", font); f.size = Pt(o.get("size", size))
        f.bold = o.get("bold", bold); f.italic = o.get("italic", False)
        f.color.rgb = o.get("color", color)

def table(slide, l, t, w, header, rows, col_w=None, size=12, header_fill=ACCENT,
          header_color=WHITE, hi=None, hi_fill=HIROW, aligns=None, font=FT,
          row_h=0.34, header_h=0.36, body_color=INK):
    """rows = list of list-of(cell runs). hi = set of row indices to highlight."""
    nC = len(header); nR = len(rows) + 1
    h = Inches(header_h + row_h * len(rows))
    gt = slide.shapes.add_table(nR, nC, l, t, w, h).table
    gt.first_row = False; gt.horz_banding = False
    # kill default style banding
    tblPr = gt._tbl.tblPr
    tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')
    if col_w:
        tot = sum(col_w)
        for j, cw in enumerate(col_w):
            gt.columns[j].width = Emu(int(w * cw / tot))
    gt.rows[0].height = Inches(header_h)
    if aligns is None:
        aligns = [PP_ALIGN.LEFT] + [PP_ALIGN.CENTER]*(nC-1)
    for j, htxt in enumerate(header):
        _cell(gt.cell(0, j), htxt, size, header_color, True, aligns[j], header_fill, font=font)
    hi = hi or set()
    for i, row in enumerate(rows):
        gt.rows[i+1].height = Inches(row_h)
        base = hi_fill if i in hi else (ROW2 if i % 2 else ROW1)
        for j, c in enumerate(row):
            _cell(gt.cell(i+1, j), c, size, body_color, False, aligns[j], base, font=font)
    return gt

def chip(slide, l, t, w, text, fill, tcolor=WHITE, size=12.5, h=0.38):
    sp = rrect(slide, l, t, w, Inches(h), fill)
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _runs(p, text if isinstance(text, list) else [(text, {})], size, tcolor, font=FT)
    return sp

# convenience for "winner ➜ pair" callouts
def winbar(slide, l, t, w, label, value, color=ACC2):
    return chip(slide, l, t, w, [(label+"  ", {"bold": True}), (value, {})], color, size=12.5)


LX = Inches(0.55); CW = Inches(12.23)
def f2(v): return f"{v:.2f}"
def f3(v): return f"{v:.3f}"

# =========================================================
# 1 — TITLE
# =========================================================
s = prs.slides.add_slide(BLANK); _page[0] += 1
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, Inches(4.55), SW, Inches(0.05), ACCENT)
textbox(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.6), [
    [("Interpolating Runtime ", {"size": 40, "bold": True, "color": WHITE}),
     ("Q-errors", {"size": 40, "bold": True, "color": RGBColor(0x6F, 0xD3, 0xDA)})],
    [("Sampling a parameter grid cheaply to estimate where the optimizer mis-predicts", {"size": 18, "color": RGBColor(0xC8, 0xD2, 0xDA)})],
], space_after=10)
textbox(s, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.6), [
    [("Q-error of a pair:   ", {"size": 17, "color": RGBColor(0x9F, 0xB0, 0xBC)}),
     ("q = max(rtₓ, rtₓ₋₁) / min(rtₓ, rtₓ₋₁)  ≥ 1", {"size": 18, "bold": True, "color": WHITE, "font": FM})],
    [("11 sampling methods · worked 5×5 grid (40 pairs) · every evaluation metric · Q5 results: parallel ON vs OFF",
      {"size": 14, "color": RGBColor(0x9F, 0xB0, 0xBC)})],
], space_after=12)
textbox(s, Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
        [[("scripts/interpolation.py   +   scripts/evaluate_interpolation.py", {"size": 12, "color": ACCENT, "font": FM})]])

# =========================================================
# 2 — THE LABEL Y (q-error)
# =========================================================
s = slide("Foundations · the label", "What we predict:  the per-pair Q-error  Y")
bullets(s, LX, Inches(1.5), Inches(6.0), Inches(3.6), [
    (0, [("An instance = an ", {}), ("axis-neighbour pair", {"bold": True}), (" of grid points (a, b).", {})]),
    (1, [("Each grid point has a measured runtime rt.", {})]),
    (0, [("Label  Y  = symmetric runtime ratio of the pair:", {})]),
    (1, [("≥ 1 always; ", {}), ("1 = identical", {"bold": True, "color": GREEN}), (" runtimes, large = a ", {}),
         ("cliff", {"bold": True, "color": ACC2}), (" (e.g. plan flip).", {})]),
    (0, [("Symmetric ⇒ direction-free; we model ", {}), ("log Y", {"bold": True}), (" (ratios are multiplicative).", {})]),
    (0, [("Goal: estimate Y for ", {}), ("all", {"italic": True}), (" pairs after measuring only a few.", {})]),
])
mono(s, Inches(6.8), Inches(1.5), Inches(5.95), Inches(1.15), [
    "q(a,b) = max(rtₐ, rt_b) / min(rtₐ, rt_b)   ≥ 1",
    "model space:  z = log q          q̂ = exp(ẑ) ≥ 1",
], size=14.5, title="symmetric ratio  (interpolation.py · symmetric_ratio)")
table(s, Inches(6.8), Inches(2.95), Inches(5.95),
      ["pair (toy)", "rtₐ", "rt_b", "q = max/min"],
      [
       [[(T.pair_label(0), {"font": FM})], "10.0", "11.2", [("1.12", {"bold": True})]],
       [[(T.pair_label(8), {"font": FM})], "13.2", "22.0", [("1.67", {"bold": True})]],
       [[(T.pair_label(22), {"font": FM})], "14.5", "110", [("7.59", {"bold": True, "color": ACC2})]],
       [[(T.pair_label(30), {"font": FM})], "14.9", "110", [("7.38", {"bold": True, "color": ACC2})]],
      ], col_w=[3.2,1,1,1.6], size=12.5)
textbox(s, Inches(6.8), Inches(4.95), Inches(5.95), Inches(0.4),
        [[("Most pairs ≈ 1 (smooth); a few pairs are large (the interesting ones).", {"size": 11.5, "italic": True, "color": GREY})]])

# =========================================================
# 3 — GRID -> PAIRS  (counting)
# =========================================================
s = slide("Foundations · the instances", "From a 5×5 grid to 40 directed-axis pairs")
mono(s, LX, Inches(1.5), Inches(6.0), Inches(1.65), [
    "points          = N × N            = 5 × 5 = 25",
    "axes            = 2  (x1, x2)",
    "pairs / axis    = (N−1) × N        = 4 × 5 = 20",
    "TOTAL pairs  P  = 2·(N−1)·N        = 40",
], size=14, title="counting the pairs  ( 4×4 'resolution' = 4 gaps per row )")
bullets(s, LX, Inches(3.35), Inches(6.0), Inches(3.2), [
    (0, [("axis 0", {"bold": True, "color": ACCENT}), ("  pairs step in x1:  (i,j)→(i+1,j)    — 20 pairs", {})]),
    (0, [("axis 1", {"bold": True, "color": ACCENT}), ("  pairs step in x2:  (i,j)→(i,j+1)    — 20 pairs", {})]),
    (0, [("Built exactly as interpolation.py: for each point, look ", {}), ("+1 along each axis", {"bold": True}), (".", {})]),
    (0, [("Same construction scales: the real Q5 grid → ", {}),
         ("P = 220 pairs", {"bold": True, "color": ACC2}), (".", {})]),
])
# mini 5x5 lattice picture (right)
gx0, gy0, step = Inches(7.55), Inches(1.7), Inches(0.92)
for i in range(5):
    for j in range(5):
        cx = gx0 + step*j; cy = gy0 + step*i
        d = rrect(s, cx, cy, Inches(0.2), Inches(0.2), INK);
        if j < 4:
            rect(s, cx+Inches(0.2), cy+Inches(0.07), step-Inches(0.2), Pt(1.4), ACCENT)
        if i < 4:
            rect(s, cx+Inches(0.07), cy+Inches(0.2), Pt(1.4), step-Inches(0.2), RGBColor(0x6F, 0xB7, 0xBD))
textbox(s, gx0, gy0+step*5+Inches(0.02), Inches(5), Inches(0.5), [
    [("●", {"color": INK, "size": 12}), (" grid point   ", {"size": 11, "color": GREY}),
     ("—", {"color": ACCENT, "size": 12, "bold": True}), (" axis-0 pair   ", {"size": 11, "color": GREY}),
     ("|", {"color": RGBColor(0x6F, 0xB7, 0xBD), "size": 12, "bold": True}), (" axis-1 pair", {"size": 11, "color": GREY})]])

# =========================================================
# 4 — FEATURES X + BUDGET
# =========================================================
s = slide("Foundations · features & budget", "Feature X per pair, and the sampling budget")
mono(s, LX, Inches(1.5), Inches(6.0), Inches(1.5), [
    "X = [ mid_x1 , mid_x2 ,  1[axis=0] , 1[axis=1] ]",
    "    └─ midpoint ─┘     └── one-hot axis ──┘",
    "fdim = 2 (midpoint) + 2 (axes) = 4",
], size=13.5, title="feature vector  ( what the GP / interpolator sees )")
mono(s, LX, Inches(3.15), Inches(6.0), Inches(1.5), [
    "MIN = max(5, fdim+2)           = max(5,6) = 6",
    "B   = max( MIN , ⌈0.10 · P⌉ )",
    "toy:  B = max(6, ⌈4⌉) = 6   →  we use 10*",
], size=13.5, title="budget  ( interpolation.py )")
textbox(s, LX, Inches(4.75), Inches(6.0), Inches(0.6),
        [[("*toy budget raised to 10/40 so the active loop is visible; production = 10% of P.",
           {"size": 11, "italic": True, "color": GREY})]])
table(s, Inches(6.85), Inches(1.5), Inches(5.9),
      ["pid", "axis", "X = [mid₁,mid₂,a₀,a₁]", "Y=q"],
      [
       ["0",  "0", [("[1.5, 1.0, 1, 0]", {"font": FM})], [("1.12", {"bold": True})]],
       ["4",  "0", [("[1.5, 3.0, 1, 0]", {"font": FM})], [("1.12", {"bold": True})]],
       ["8",  "0", [("[1.5, 5.0, 1, 0]", {"font": FM})], [("1.67", {"bold": True})]],
       ["22", "0", [("[3.5, 3.0, 1, 0]", {"font": FM})], [("7.59", {"bold": True, "color": ACC2})]],
       ["30", "1", [("[4.0, 2.5, 0, 1]", {"font": FM})], [("7.38", {"bold": True, "color": ACC2})]],
      ], col_w=[0.7,0.8,3.2,1.0], size=12.5)
textbox(s, Inches(6.85), Inches(3.62), Inches(5.9), Inches(0.3),
        [[("master pair table  ·  5 of P=40 rows", {"size": 11.5, "italic": True, "color": GREY})]])
bullets(s, Inches(6.85), Inches(4.05), Inches(5.9), Inches(2.4), [
    (0, [("All methods see ", {}), ("only", {"bold": True, "color": RED}), (" the Y of pairs they sample.", {})]),
    (0, "Two ways to turn samples → full prediction:"),
    (1, [("interpolate", {"bold": True}), (" between samples (random, stride), or", {})]),
    (1, [("fit a model", {"bold": True}), (" + actively choose samples (GP / BO).", {})]),
])

# =========================================================
# 4b — WHAT ELSE COULD GO IN THE FEATURE VECTOR
# =========================================================
s = slide("Foundations · richer features", "What else could go into the feature X?")
bullets(s, LX, Inches(1.45), Inches(12.23), Inches(0.95), [
    (0, [("Today ", {}), ("X = midpoint + axis one-hot", {"font": FM, "bold": True}),
         (" — pure geometry.  Each pair's two endpoints already carry much more signal in ", {}),
         ("ground_truth.csv", {"font": FM}), (" — adding it lets the GP ", {}),
         ("smell a cliff before sampling", {"bold": True}), (".", {})]),
])
table(s, LX, Inches(2.5), Inches(12.23),
      ["signal group", "columns (per endpoint, + Δ across the pair)", "why it predicts the pair's q-error"],
      [
       [[("Selectivity gap", {"bold": True})],
        [("selectivity_x*, joint_sel, dS_*, joint_dS", {"font": FM, "size": 10.5})],
        "a big cardinality jump across the pair → runtime jump"],
       [[("Optimizer estimate", {"bold": True})],
        [("plan_rows, total_cost, startup_cost (+ ratio)", {"font": FM, "size": 10.5})],
        "where the optimizer's own guess swings, q-error spikes"],
       [[("Plan structure", {"bold": True, "color": ACC2})],
        [("plan_change, node_count, max_depth, join/scan/hash/sort", {"font": FM, "size": 10})],
        [("a plan flip across the pair is the #1 cliff signal", {"bold": True})]],
       [[("Parallelism / IO", {"bold": True})],
        [("workers_planned/launched, shared_hit_ratio, temp_total_blocks", {"font": FM, "size": 10})],
        "parallel kick-in or a spill changes runtime sharply"],
       [[("Pair geometry", {"bold": True})],
        [("raw value-gap width, boundary-vs-interior, axis (have)", {"font": FM, "size": 10.5})],
        "wide gaps and grid boundaries are higher-risk"],
      ], col_w=[2.1, 5.0, 5.1], size=11.5, row_h=0.6,
      aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT])
textbox(s, LX, Inches(6.2), Inches(12.23), Inches(0.85), [
    [("Rule:", {"bold": True, "color": RED}), (" only features knowable ", {}),
     ("without running the pair", {"bold": True}), (" — EXPLAIN-derived (plan, cost, plan_rows) and cached selectivity are cheap & legal; the pair's ", {}),
     ("actual runtime / q-error must never leak.", {"bold": True}),
     ("   Standardize numerics (log for rows/cost), one-hot categoricals (root_node).", {"color": GREY})],
])

# =========================================================
# 5 — EVALUATION METRICS  I  (accuracy)
# =========================================================
s = slide("Evaluation · 1 of 2", "How every method is scored — prediction accuracy", accent=ACC2)
mono(s, LX, Inches(1.5), Inches(5.7), Inches(1.15), [
    "pred q-error:  Q = max( y/ŷ , ŷ/y )  ≥ 1",
    "ALL  = every pair      UNSEEN = only un-sampled pairs",
], size=13.5, title="the headline number  (evaluate_interpolation.py)", accent=ACC2)
bullets(s, LX, Inches(2.85), Inches(5.7), Inches(3.8), [
    (0, [("UNSEEN", {"bold": True, "color": ACC2}), (" is the real test — accuracy where we ", {}),
         ("did not", {"italic": True}), (" measure.", {})]),
    (0, [("Ranking column = ", {}), ("QERR_MEDIAN_UNSEEN", {"bold": True, "font": FM, "size": 13})]),
    (1, "lower = predicted q-errors are closer to truth."),
    (0, "WITHIN_kX = % of pairs whose prediction is within k×."),
])
table(s, Inches(6.5), Inches(1.5), Inches(6.25),
      ["metric", "formula", "best"],
      [
       ["MAE",          [("mean |y − ŷ|", {"font": FM})], "0"],
       ["RMSE",         [("√ mean(y − ŷ)²", {"font": FM})], "0"],
       ["R²",           [("1 − SS_res/SS_tot", {"font": FM})], "1"],
       ["MAPE",         [("mean |y−ŷ|/y · 100", {"font": FM})], "0"],
       ["QERR_MEAN/MED", [("mean / median of Q", {"font": FM})], "1"],
       ["QERR_P90/P95/MAX", [("tail percentiles of Q", {"font": FM})], "1"],
       ["WITHIN_2X/5X", [("% ( Q ≤ 2 ) , % ( Q ≤ 5 )", {"font": FM})], "100"],
      ], col_w=[2.1,3.6,0.8], size=12, header_fill=ACC2)

# =========================================================
# 6 — EVALUATION METRICS II  (discovery + smoothness)
# =========================================================
s = slide("Evaluation · 2 of 2", "Discovery of bad pairs  &  global smoothness summaries", accent=ACC2)
table(s, LX, Inches(1.5), Inches(6.05),
      ["discovery metric", "meaning"],
      [
       ["max_qerr_found", [("largest true q among ", {}), ("sampled", {"bold": True}), (" pairs", {})]],
       ["coverage_2x", [("% of pairs with q>2 that were sampled", {})]],
       ["coverage_5x", [("% of pairs with q>5 that were sampled", {})]],
      ], col_w=[2.2,4.0], size=12, header_fill=ACC2, row_h=0.5)
mono(s, LX, Inches(3.55), Inches(6.05), Inches(1.7), [
    "S_max  = max  q           (worst pair)",
    "S_avg  = mean q           (overall roughness)",
    "S_topk = mean top-k q     (k = 10% of P, the tail)",
    "RelErr = | est − true | / true · 100%",
], size=13, title="smoothness summaries  S  (estimate vs truth)", accent=ACC2)
bullets(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(5.3), [
    (0, [("Two different jobs, scored separately:", {"bold": True})]),
    (1, [("interpolation", {"bold": True, "color": ACCENT}), (" — be accurate on every pair (UNSEEN Q).", {})]),
    (1, [("discovery / S", {"bold": True, "color": ACC2}), (" — find & quantify the worst pairs.", {})]),
    (0, "A method can win one and lose the other:"),
    (2, "pure-σ explores → great coverage, mediocre Q."),
    (2, "interpolators → smooth Q, miss the cliffs."),
    (0, [("The closing Q5 slides report ", {}), ("both", {"bold": True}), (" families of metric.", {})]),
])

# =========================================================
# 7 — SHARED MACHINERY: two predictors
# =========================================================
s = slide("Shared machinery · 1 of 2", "Two predictors underneath every method")
mono(s, LX, Inches(1.5), Inches(6.0), Inches(1.7), [
    "log q̂(X) = Σᵢ wᵢ · log qᵢ      (barycentric weights)",
    "  · linear inside the convex hull of samples",
    "  · nearest-neighbour fallback outside",
    "q̂ = exp( log q̂ ) ≥ 1",
], size=13, title="A.  Linear interpolation  →  random, uniform_stride")
mono(s, LX, Inches(3.4), Inches(6.0), Inches(3.0), [
    "k(r) = (1 + √5·r/ℓ + 5r²/3ℓ²) · exp(−√5·r/ℓ)",
    "        r = ‖Xᵢ − Xⱼ‖   (Matérn 5/2)",
    "",
    "μ(X) = k*ᵀ K⁻¹ y                 (posterior mean)",
    "σ²(X)= 1 + noise − k*ᵀ K⁻¹ k*    (posterior var)",
    "",
    "(ℓ, noise) by MLE · fit on log q · q̂ = exp(μ)",
], size=12.5, title="B.  Matérn-5/2 Gaussian Process  →  gpr & all BO")
bullets(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(5.0), [
    (0, [("Interpolation", {"bold": True, "color": ACCENT}), (" = cheap, no uncertainty.", {})]),
    (1, "good when samples are dense & spread."),
    (0, [("GP", {"bold": True, "color": ACCENT}), (" gives μ (best guess) ", {}),
         ("and", {"italic": True}), (" σ (uncertainty).", {})]),
    (1, [("σ", {"bold": True}), (" is what lets a method ", {}), ("choose", {"italic": True}), (" the next pair.", {})]),
    (1, "everything is done in log-q space, then exp()."),
    (0, [("Why log?", {"bold": True}), (" q-errors are ratios; log makes 2× and ½× symmetric.", {})]),
    (0, [("Data-starved caveat:", {"bold": True, "color": ACC2}),
         (" on tiny seeds the MLE length-scale collapses → flat μ/σ. Real grids (P=220) behave well.", {})]),
])

# =========================================================
# 7b — BO internals: the Matérn kernel family
# =========================================================
s = slide("BO internals · kernels", "The kernel: how similar are two pairs?")
bullets(s, LX, Inches(1.45), Inches(12.23), Inches(1.25), [
    (0, [("A kernel ", {}), ("k(r)", {"font": FM, "bold": True}),
         (" sets how strongly the q-errors at two pairs co-vary, given feature distance ", {}),
         ("r = ‖X − X'‖", {"font": FM}), (".  It decides how ", {}),
         ("smooth", {"bold": True}), (" the predicted surface is.", {})]),
    (0, [("The ", {}), ("Matérn", {"bold": True}), (" family is indexed by ν — larger ν ⇒ smoother paths.  ", {}),
         ("ℓ", {"font": FM, "bold": True}), (" = length-scale (how far one sample's influence reaches).", {})]),
])
table(s, LX, Inches(2.85), Inches(12.23),
      ["kernel (ν)", "k(r)", "smoothness", "used in this code"],
      [
       [[("Matérn ½", {"bold": True})], [("exp(−r/ℓ)", {"font": FM})],
        "continuous, non-diff (rough)", [("OtterTune kernel", {"font": FM})]],
       [[("Matérn 3/2", {"bold": True})], [("(1 + √3·r/ℓ)·exp(−√3·r/ℓ)", {"font": FM})],
        "1× differentiable", "—"],
       [[("Matérn 5/2", {"bold": True, "color": ACC2})],
        [("(1 + √5·r/ℓ + 5r²/3ℓ²)·exp(−√5·r/ℓ)", {"font": FM})],
        "2× differentiable", [("MaternGP — bo_*, smoothness_*, gpr_fixed", {"font": FM, "size": 11})]],
       [[("RBF / SE  (ν→∞)", {"bold": True})], [("exp(−r²/2ℓ²)", {"font": FM})],
        "∞ differentiable (very smooth)", "—"],
      ], col_w=[1.7, 4.5, 2.6, 3.4], size=12, row_h=0.62, hi={2})
textbox(s, LX, Inches(5.95), Inches(12.23), Inches(1.0), [
    [("All Matérn-based methods here use ", {"size": 13}),
     ("ν = 5/2", {"font": FM, "bold": True, "size": 13, "color": ACC2}),
     (" — smooth enough for runtime surfaces, sharp enough to bend around a cliff.  ", {"size": 13}),
     ("OtterTune uses ν = ½", {"font": FM, "bold": True, "size": 13}),
     (" (the plain exponential kernel).", {"size": 13})],
    [("Why not RBF?  Its infinite smoothness over-smooths the sharp q-error cliffs we care about.",
      {"size": 12, "italic": True, "color": GREY})],
], space_after=6)

# =========================================================
# 7c — BO internals: Matérn 5/2 in depth + OtterTune contrast
# =========================================================
s = slide("BO internals · kernels", "Matérn-5/2 — the workhorse, and OtterTune's ν=½")
mono(s, LX, Inches(1.5), Inches(6.0), Inches(1.7), [
    "a = √5 · r / ℓ        r = ‖X − X'‖  (z-scored X)",
    "",
    "k(r) = ( 1 + a + a²/3 ) · e^(−a)",
], size=14, title="Matérn 5/2  (MaternGP._k52)")
bullets(s, LX, Inches(3.35), Inches(6.0), Inches(3.4), [
    (0, [("1", {"font": FM, "bold": True}), (" baseline correlation · ", {}),
         ("a", {"font": FM, "bold": True}), (" linear · ", {}),
         ("a²/3", {"font": FM, "bold": True}), (" quadratic term.", {})]),
    (0, [("Two derivatives ⇒ models ", {}), ("curvature", {"bold": True}),
         (" but stays finite at a cliff — ideal for q-error surfaces.", {})]),
    (0, [("ℓ", {"font": FM, "bold": True}), (" learned per-fit; ", {}),
         ("r", {"font": FM}), (" in standardized feature space so all axes are comparable.", {})]),
    (0, [("Signal scaled to 1 (log-q is z-scored: ymean, ystd); ", {}),
         ("noise·I", {"font": FM}), (" added on the diagonal.", {})]),
])
mono(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(1.7), [
    "K = mag · e^(−‖ΔX‖/ℓ) + ridge · I",
    "",
    "ℓ=2,  mag=1,  ridge=1   (fixed, NO MLE)",
], size=14, title="OtterTune  ν=½  (exponential)")
bullets(s, Inches(6.85), Inches(3.35), Inches(5.9), Inches(3.4), [
    (0, [("Same shape as Matérn ½: ", {}), ("e^(−r/ℓ)", {"font": FM}), (".", {})]),
    (0, [("Non-differentiable ⇒ ", {}), ("rougher", {"bold": True}),
         (" paths; reacts fast to local jumps.", {})]),
    (0, [("Hyper-parameters are ", {}), ("fixed", {"bold": True, "color": ACC2}),
         (" (ported from cmu-db/ottertune), not fit by MLE.", {})]),
    (0, [("ridge·I", {"font": FM}), (" is the noise/jitter that keeps K invertible.", {})]),
])

# =========================================================
# 7d — BO internals: MLE fit + posterior
# =========================================================
s = slide("BO internals · fitting", "Choosing (ℓ, noise) by marginal likelihood, then predicting")
mono(s, LX, Inches(1.5), Inches(6.0), Inches(2.5), [
    "K = k(D; ℓ) + noise · I        (D = pair-distances)",
    "",
    "log p(y|X) = −½ yᵀK⁻¹y           (data fit)",
    "             − ½ ln|K|            (complexity)",
    "             − (n/2) ln 2π",
], size=13, title="marginal log-likelihood  (MaternGP._lml)")
bullets(s, LX, Inches(4.2), Inches(6.0), Inches(2.5), [
    (0, [("Grid-search ", {}), ("ℓ ∈ logspace(−1, 1.3, 12)", {"font": FM, "size": 12}),
         (",  ", {}), ("noise ∈ {1e−6 … 1e−1}", {"font": FM, "size": 12}), (".", {})]),
    (0, "Keep the (ℓ, noise) with the highest log-likelihood."),
    (0, [("Cholesky factor of K for a stable solve (", {}),
         ("yᵀK⁻¹y", {"font": FM}), (", ln|K|).", {})]),
])
mono(s, Inches(6.85), Inches(1.5), Inches(5.9), Inches(2.5), [
    "μ(X)  = k*ᵀ K⁻¹ y",
    "σ²(X) = k(X,X) + noise − k*ᵀ K⁻¹ k*",
    "",
    "back-transform:",
    "  q̂ = exp( μ·ystd + ymean )  ≥ 1",
], size=13, title="posterior at a new pair X  (predict)")
bullets(s, Inches(6.85), Inches(4.2), Inches(5.9), Inches(2.5), [
    (0, [("k*", {"font": FM, "bold": True}), (" = kernel from X to every sampled pair.", {})]),
    (0, [("μ", {"font": FM, "bold": True}), (" = best guess (log-q); ", {}),
         ("σ", {"font": FM, "bold": True}), (" = uncertainty.", {})]),
    (0, [("This ", {}), ("(μ, σ)", {"font": FM, "bold": True}),
         (" is the only input every acquisition needs →", {})]),
])

# =========================================================
# 8 — SHARED MACHINERY: the active loop
# =========================================================
s = slide("Shared machinery · 2 of 2", "The active-sampling loop — only the acquisition changes")
mono(s, LX, Inches(1.5), Inches(7.35), Inches(2.5), [
    "seed  ← LHS-maximin pairs        (space-filling, n=6)",
    "while |sampled| < BUDGET:",
    "    fit GP on sampled (X, log q)",
    "    μ, σ ← GP.predict(all pairs)",
    "    score ← ACQUISITION(μ, σ)   ◄ the ONE differing line",
    "    next  ← argmax score (un-sampled)",
    "    measure next ; add to sampled",
], size=13, title="generic active loop  (bo_generic / smoothness_* / ottertune)")
bullets(s, Inches(8.15), Inches(1.6), Inches(4.65), Inches(2.4), [
    (0, [("LHS-maximin", {"bold": True}), (": a spread-out start so the GP isn't blind anywhere.", {})]),
    (0, [("Passive", {"bold": True}), (" methods skip the loop — fixed seed, predict once.", {})]),
    (0, [("Greedy", {"bold": True}), (": one pair per iteration, then refit.", {})]),
])
textbox(s, LX, Inches(4.25), Inches(11), Inches(0.3),
        [[("Acquisitions covered in this deck", {"size": 13.5, "bold": True, "color": INK}),
          ("   (the only thing that changes):", {"size": 13.5, "color": GREY})]])
chips = [
    ("σ", "explore most-uncertain", ACCENT),
    ("UCB", "μ + κσ  exploit + explore", ACCENT),
    ("EI", "expected gain vs best-so-far", ACCENT),
    ("β-UCB", "growing β  (smoothness_max)", ACC2),
    ("BQ", "represent the mean  (avg)", ACC2),
    ("EI-τ", "tail threshold  (topk)", ACC2),
]
cw = Inches(3.95)
for i, (a, b, c) in enumerate(chips):
    col = i % 3; rowi = i // 3
    x = LX + col * (cw + Inches(0.19))
    yy = Inches(4.7) + rowi * Inches(0.62)
    chip(s, x, yy, cw, [(a + ":  ", {"bold": True}), (b, {})], c, size=12, h=0.5)

# =========================================================
# content helpers for the method slides
# =========================================================
def Xstr(pid):
    m = T.mid[pid]; a0 = 1 if T.axis_arr[pid] == 0 else 0
    return f"[{m[0]:.1f},{m[1]:.1f},{a0},{1-a0}]"

def qcell(q, bold=True):
    return [(f"{q:.2f}", {"bold": bold, "color": ACC2 if q > 2 else INK})]

def seed_table(s, l, t, w, ids, size=11.5, nmax=6):
    rows = []
    for pid in ids[:nmax]:
        rows.append([str(pid), [(T.pair_label(pid), {"font": FM})],
                     [(Xstr(pid), {"font": FM})], qcell(T.qerr_all[pid])])
    if len(ids) > nmax:
        rows.append(["…", [(f"+{len(ids)-nmax} more", {"italic": True, "color": GREY})], "", ""])
    return table(s, l, t, w, ["pid", "pair", "X = [mid₁,mid₂,a₀,a₁]", "Y=q"],
                 rows, col_w=[0.6, 2.6, 2.5, 0.8], size=size, row_h=0.3)

def pred_table(s, l, t, w, ids_show, pred, size=12, row_h=0.32):
    rows = []
    for pid in ids_show:
        yt = T.qerr_all[pid]; yp = pred[pid]
        Q = max(yt/yp, yp/yt)
        rows.append([[(T.pair_label(pid), {"font": FM})], f"{yt:.2f}", f"{yp:.2f}",
                     [(f"{Q:.2f}", {"bold": True, "color": RED if Q > 2 else GREEN})]])
    return table(s, l, t, w, ["unsampled pair", "y", "ŷ", "Q=max(y/ŷ,ŷ/y)"],
                 rows, col_w=[2.6, 0.8, 0.8, 1.7], size=size, row_h=row_h)

def gp_pred_table(s, l, t, w, ids_show, res, size=12, row_h=0.32):
    rows = []
    for pid in ids_show:
        yt = T.qerr_all[pid]; yp = res["pred"][pid]; sd = res["sd"][pid]
        rows.append([[(T.pair_label(pid), {"font": FM})], f"{yt:.2f}",
                     f"{yp:.2f}", f"{sd:.2f}"])
    return table(s, l, t, w, ["unsampled pair", "y", "q̂=eᵘ", "σ"],
                 rows, col_w=[2.6, 0.8, 1.0, 0.9], size=size, row_h=row_h)

def trace_table(s, l, t, w, res, size=11.5):
    seed = res["seed"]; picks = res["picks"]; B = T.BUDGET
    rows = [["0", [("LHS seed ×%d" % len(seed), {"italic": True})], "—",
             f"{len(seed)}/{B}", str(B-len(seed))]]
    hi = set()
    for k, pid in enumerate(picks, 1):
        nso = len(seed)+k; q = T.qerr_all[pid]
        if q > 2: hi.add(k)
        rows.append([str(k), [(T.pair_label(pid), {"font": FM})], qcell(q),
                     f"{nso}/{B}", str(B-nso)])
    return table(s, l, t, w, ["it", "picked pair", "q", "used", "left"],
                 rows, col_w=[0.45, 2.5, 0.65, 0.8, 0.6], size=size, row_h=0.3, hi=hi)

# shared illustrative acquisition table (μ,σ + chosen columns)
def acq_table(s, l, t, w, cols, winner, size=12):
    A = T.ACQ
    header = ["", "pair", "q̂", "μ"] + [c[0] for c in cols]
    rows = []
    for i, (tag, pid, lbl, qt) in enumerate(T.ACQ_CAND):
        row = [[(tag, {"bold": True, "color": ACCENT})],
               [(lbl, {"font": FM, "size": size-1.5})],
               f"{A['qhat'][i]:.2f}", f"{A['mu'][i]:.2f}"]
        for (h, key, fmt) in cols:
            row.append([(fmt(A[key][i]), {"bold": i == winner})])
        rows.append(row)
    colw = [0.4, 2.3, 0.6, 0.6] + [1.0]*len(cols)
    return table(s, l, t, w, header, rows, col_w=colw, size=size, row_h=0.33,
                 hi={winner})

ACQ_NOTE = "μ,σ = representative GP posterior (log-q) over 5 toy pairs; exact acquisition formulas."

# =========================================================
# 8a — BO acquisition functions, all formulas in one place
# =========================================================
s = slide("BO internals · acquisitions", "Every acquisition function in one place")
table(s, LX, Inches(1.4), Inches(12.23),
      ["acquisition", "formula  (maximise, unless noted)", "goal", "method(s)"],
      [
       [[("σ  uncertainty", {"bold": True})], [("a = σ(X)", {"font": FM})],
        "explore", [("bo_sigma, ottertune_reconstruct", {"font": FM, "size": 10.5})]],
       [[("UCB", {"bold": True})], [("a = μ + κ·σ        κ = 2 (fixed)", {"font": FM})],
        "balance", [("bo_generic_ucb", {"font": FM, "size": 10.5})]],
       [[("EI", {"bold": True})], [("a = (μ−f⁺)Φ(z) + σφ(z),  z=(μ−f⁺)/σ", {"font": FM, "size": 11})],
        "beat best f⁺", [("bo_generic_ei", {"font": FM, "size": 10.5})]],
       [[("GP-UCB  (adaptive)", {"bold": True, "color": ACC2})],
        [("a = μ + β_t·σ,  β_t=√(2 ln(P t²π²/6))", {"font": FM, "size": 11})],
        "S_max (peak)", [("smoothness_max", {"font": FM, "size": 10.5})]],
       [[("Bayesian Quadrature", {"bold": True, "color": ACC2})],
        [("a = (Σ_q k*(q,p))² / σ*²(p)", {"font": FM})],
        "S_avg (mean)", [("smoothness_avg", {"font": FM, "size": 10.5})]],
       [[("EI-threshold", {"bold": True, "color": ACC2})],
        [("a = (μ−τ_t)Φ(z)+σφ(z),  τ_t = k-th largest", {"font": FM, "size": 10.5})],
        "S_topk (tail)", [("smoothness_topk", {"font": FM, "size": 10.5})]],
       [[("OtterTune GPRGD", {"bold": True})],
        [("minimise μ_mult·μ − β·σ  over continuous x", {"font": FM, "size": 11})],
        "max / cover", [("ottertune_max / reconstruct", {"font": FM, "size": 10.5})]],
      ], col_w=[2.5, 5.3, 1.7, 3.4], size=12, row_h=0.56,
      aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT])
textbox(s, LX, Inches(5.95), Inches(12.23), Inches(1.0), [
    [("Φ", {"font": FM, "bold": True}), (" = normal CDF   ", {}),
     ("φ", {"font": FM, "bold": True}), (" = normal PDF   ", {}),
     ("f⁺", {"font": FM, "bold": True}), (" = best q seen   ", {}),
     ("τ_t", {"font": FM, "bold": True}), (" = tail threshold   ", {}),
     ("t", {"font": FM, "bold": True}), (" = iteration   ", {}),
     ("P", {"font": FM, "bold": True}), (" = #pairs", {})],
    [("Same GP posterior (μ, σ) feeds all of them — only this scoring rule changes which pair is sampled next.",
      {"italic": True, "color": GREY})],
], space_after=6)

# =========================================================
# 9 — METHODS MAP
# =========================================================
s = slide("Methods · map", "Eleven methods in four families")
table(s, LX, Inches(1.5), Inches(12.23),
      ["family", "methods", "chooses\nsamples?", "predictor / idea"],
      [
       [[("Passive baselines", {"bold": True})], [("random , uniform_stride", {"font": FM})],
        [("no", {"color": RED, "bold": True})], "linear interpolation between fixed samples"],
       [[("Passive GP", {"bold": True})], [("gpr_fixed", {"font": FM})],
        [("no", {"color": RED, "bold": True})], "one-shot Matérn GP on a uniform seed"],
       [[("Active · generic BO", {"bold": True, "color": ACCENT})],
        [("bo_generic_sigma / ucb / ei", {"font": FM})],
        [("yes", {"color": GREEN, "bold": True})], "Matérn GP + acquisition (σ / UCB / EI)"],
       [[("Active · OtterTune", {"bold": True, "color": ACCENT})],
        [("ottertune_max / reconstruct", {"font": FM})],
        [("yes", {"color": GREEN, "bold": True})], "fixed-HP kernel + gradient search (GPRGD)"],
       [[("Active · smoothness", {"bold": True, "color": ACCENT})],
        [("smoothness_max / avg / topk", {"font": FM})],
        [("yes", {"color": GREEN, "bold": True})], "GP + acquisition tuned to S_max / S_avg / S_topk"],
      ], col_w=[2.3, 3.4, 1.1, 4.6], size=12.5, row_h=0.62)
textbox(s, LX, Inches(5.7), Inches(12.23), Inches(0.9), [
    [("ground_truth", {"bold": True, "font": FM, "size": 13}),
     (" measures all P pairs — the baseline every method is scored against.", {"size": 13})],
    [("Passive = predict once from a fixed sample set.   Active = repeat {fit → score → sample} until the budget is spent.",
      {"size": 13, "color": GREY, "italic": True})],
], space_after=6)

# =========================================================
# 10 — random
# =========================================================
s = slide("Baseline 1 · random  (passive)", "random — blind sample, then interpolate")
bullets(s, LX, Inches(1.5), Inches(5.9), Inches(2.2), [
    (0, [("Draw ", {}), ("B = 10", {"bold": True, "font": FM}), (" pair-ids uniformly at random.", {})]),
    (0, "Measure their q-errors."),
    (0, "Predict every pair by linear interpolation of log q."),
])
mono(s, LX, Inches(3.7), Inches(5.9), Inches(1.0), [
    "log q̂(X) = Σᵢ wᵢ · log qᵢ        q̂ = exp(·)",
], size=13, title="predictor — linear interpolation")
textbox(s, LX, Inches(4.85), Inches(5.9), Inches(1.7), [
    [("No targeting:", {"bold": True, "color": RED}), (" samples land mostly in the flat region;", {})],
    [("the cliff pairs are usually missed, so their q is predicted ≈ 1.", {})],
], space_after=4)
seed_table(s, Inches(6.85), Inches(1.5), Inches(5.9), T.random_ids, nmax=4)
textbox(s, Inches(6.85), Inches(3.42), Inches(5.9), Inches(0.3),
        [[("sampled (4 of 10 shown)", {"size": 11, "italic": True, "color": GREY})]])
pred_table(s, Inches(6.85), Inches(3.8), Inches(5.9), [0, 3, 22, 30, 33], T.RES["random_pred"], row_h=0.3)
textbox(s, Inches(6.85), Inches(5.74), Inches(5.9), Inches(0.7),
        [[("cliff pairs 22 & 30 (q≈7.5) predicted ≈ 1.1  →  Q ≈ 7  (a miss).",
           {"size": 11.5, "italic": True, "color": RED})]])

# =========================================================
# 11 — uniform_stride
# =========================================================
s = slide("Baseline 2 · uniform_stride  (passive)", "uniform_stride — even cover, then interpolate")
bullets(s, LX, Inches(1.5), Inches(5.9), Inches(2.3), [
    (0, [("Take ", {}), ("B", {"bold": True, "font": FM}), (" pair-ids evenly along the catalog:", {})]),
    (1, [("ids = round( linspace(0, P−1, B) )", {"font": FM, "size": 13})]),
    (0, "Same linear-interpolation predictor as random."),
    (0, [("Even in id-space ", {}), ("≠", {"bold": True}), (" even in feature-space:", {})]),
    (1, [("bracketed cliffs are nailed (24), but the cliff ", {}),
         ("smears", {"bold": True, "color": RED}), (" onto flat pairs (35, 33) and 22 is missed.", {})]),
])
mono(s, LX, Inches(4.7), Inches(5.9), Inches(0.95), [
    "ids = ⌊ linspace(0, P−1, B) ⌋   (deduped)",
], size=13, title="deterministic seed")
seed_table(s, Inches(6.85), Inches(1.5), Inches(5.9), T.stride_ids, nmax=4)
textbox(s, Inches(6.85), Inches(3.42), Inches(5.9), Inches(0.3),
        [[("sampled (4 of 10 shown) — by luck includes cliff pairs 26, 30", {"size": 11, "italic": True, "color": GREY})]])
pred_table(s, Inches(6.85), Inches(3.8), Inches(5.9), [24, 22, 35, 33, 16], T.RES["stride_pred"], row_h=0.3)
textbox(s, Inches(6.85), Inches(5.74), Inches(5.9), Inches(0.95),
        [[("24 nailed (bracketed); ", {"size": 11.5, "color": GREEN}),
          ("35 & 33 are false alarms (1.1 → 7.6); ", {"size": 11.5, "color": RED}),
          ("22 still missed.", {"size": 11.5, "color": RED})]])

# =========================================================
# 12 — gpr_fixed
# =========================================================
s = slide("Baseline 3 · gpr_fixed  (passive GP)", "gpr_fixed — one-shot Matérn GP on the stride seed")
bullets(s, LX, Inches(1.5), Inches(5.9), Inches(2.2), [
    (0, [("Seed = uniform_stride pairs (", {}), ("no", {"bold": True, "color": RED}), (" active picks).", {})]),
    (0, "Fit Matérn-5/2 GP on (X, log q); predict all pairs once."),
    (0, [("Gives μ ", {}), ("and", {"italic": True}), (" σ, but never uses σ to sample more.", {})]),
])
mono(s, LX, Inches(3.7), Inches(5.9), Inches(1.55), [
    "μ(X) = k*ᵀ K⁻¹ y          q̂ = exp(μ)",
    "σ²(X)= 1 + noise − k*ᵀ K⁻¹ k*",
], size=13, title="posterior (predict once)")
textbox(s, LX, Inches(5.35), Inches(5.9), Inches(1.55), [
    [("Data-starved toy:", {"bold": True, "color": ACC2}),
     (" 10 seeds → the MLE length-scale collapses (ℓ→0.1), so the GP returns the ",
      {}), ("global mean q̂≈1.71", {"bold": True}), (" with high σ everywhere.", {})],
    [("On the real P=220 grid it is well-conditioned (closing slides: Q̃=1.15).",
      {"italic": True, "color": GREY, "size": 12})],
], space_after=4)
gp_pred_table(s, Inches(6.85), Inches(1.5), Inches(5.9), [24, 22, 1, 35, 16], T.RES["gpr_fixed"], row_h=0.3)
textbox(s, Inches(6.85), Inches(3.45), Inches(5.9), Inches(0.75),
        [[("true q swings 1.1→7.6, but q̂ is flat at 1.71 + high σ: the GP says ",
           {"size": 11.5, "italic": True, "color": GREY}),
          ("'unsure'", {"size": 11.5, "italic": True, "bold": True}),
          (" — yet, passive, it cannot resample.", {"size": 11.5, "italic": True, "color": GREY})]])
chip(s, Inches(6.85), Inches(4.25), Inches(5.9),
     [("Active methods (next) turn that σ into the next sample.", {"bold": True})],
     ACCENT, size=12.5, h=0.5)

# =========================================================
# 13 — BO family overview (shared posterior + seed)
# =========================================================
s = slide("Active · generic BO", "bo_generic — shared loop, one decision step")
bullets(s, LX, Inches(1.5), Inches(5.85), Inches(2.0), [
    (0, [("Seed ", {}), ("6 LHS-maximin", {"bold": True}), (" pairs → fit Matérn GP.", {})]),
    (0, "Score all unsampled pairs by an acquisition of (μ, σ)."),
    (0, [("Sample argmax, refit, repeat to B=10  (", {}),
         ("σ / UCB / EI", {"bold": True, "color": ACCENT}), (" differ only here).", {})]),
])
seed_table(s, LX, Inches(3.55), Inches(5.85), T.RES["bo_sigma"]["seed"], nmax=6)
textbox(s, LX, Inches(5.78), Inches(5.85), Inches(0.3),
        [[("LHS seed for bo_generic_sigma (each variant reseeds)", {"size": 11, "italic": True, "color": GREY})]])
# shared posterior panel
table(s, Inches(6.95), Inches(1.5), Inches(5.8),
      ["", "candidate pair", "q̂=eᵘ", "μ", "σ"],
      [[[(t, {"bold": True, "color": ACCENT})], [(lbl, {"font": FM, "size": 11})],
        f"{T.ACQ['qhat'][i]:.2f}", f"{T.ACQ['mu'][i]:.2f}", f"{T.ACQ['sd'][i]:.2f}"]
       for i, (t, pid, lbl, qt) in enumerate(T.ACQ_CAND)],
      col_w=[0.4, 2.6, 0.7, 0.7, 0.7], size=12, row_h=0.34)
textbox(s, Inches(6.95), Inches(3.55), Inches(5.8), Inches(0.55),
        [[("a representative posterior — the next 3 slides apply σ / UCB / EI to ",
           {"size": 11.5, "italic": True, "color": GREY}),
          ("these same μ,σ.", {"size": 11.5, "italic": True, "bold": True, "color": INK})]])
bullets(s, Inches(6.95), Inches(4.15), Inches(5.8), Inches(2.4), [
    (1, [("A,B", {"bold": True, "color": ACC2}), (" sit on the cliff (high q): high μ.", {})]),
    (1, [("E", {"bold": True}), (" is a far, unsampled low corner: high σ, low μ.", {})]),
    (1, [("D", {"bold": True}), (" sits among samples: low μ, low σ.", {})]),
])

# =========================================================
# 14 — bo_generic_sigma
# =========================================================
s = slide("Active BO · σ", "bo_generic_sigma — sample where the GP is least sure")
mono(s, LX, Inches(1.5), Inches(7.5), Inches(0.8), [
    "acq(X) = σ(X)            next = argmax σ",
], size=15, title="acquisition")
acq_table(s, LX, Inches(2.5), Inches(7.5), [("σ = acq", "sigma", f2)], winner=4, size=12.5)
winbar(s, Inches(8.35), Inches(2.5), Inches(4.4), "picks", "E  (max σ = 0.95)")
bullets(s, Inches(8.35), Inches(3.05), Inches(4.45), Inches(1.35), [
    (1, "pure exploration: ignores μ entirely."),
    (1, [("can 'waste' budget on uncertain-but-flat corners (E).", {"color": RED})]),
    (1, "great coverage, not focused on the worst pairs."),
])
trace_table(s, Inches(8.35), Inches(4.55), Inches(4.45), T.RES["bo_sigma"])
textbox(s, LX, Inches(6.55), Inches(7.5), Inches(0.4),
        [[(ACQ_NOTE, {"size": 10.5, "italic": True, "color": GREY})]])

# =========================================================
# 15 — bo_generic_ucb
# =========================================================
s = slide("Active BO · UCB", "bo_generic_ucb — balance high guess and high doubt")
mono(s, LX, Inches(1.5), Inches(7.5), Inches(0.8), [
    "acq(X) = μ(X) + κ·σ(X)        κ = 2   (fixed)",
], size=15, title="acquisition  (upper confidence bound)")
acq_table(s, LX, Inches(2.5), Inches(7.5),
          [("σ", "sigma", f2), ("μ+κσ", "ucb", f2)], winner=1, size=12.5)
winbar(s, Inches(8.35), Inches(2.5), Inches(4.4), "picks", "B  (μ+2σ = 2.90)")
bullets(s, Inches(8.35), Inches(3.05), Inches(4.45), Inches(1.35), [
    (1, "exploit (μ) + explore (σ) in one score."),
    (1, [("prefers the ", {}), ("uncertain end of the cliff", {"bold": True}), (" (B) over flat E.", {})]),
    (1, "κ fixed → constant explore/exploit mix."),
])
trace_table(s, Inches(8.35), Inches(4.55), Inches(4.45), T.RES["bo_ucb"])
textbox(s, LX, Inches(6.55), Inches(7.5), Inches(0.4),
        [[(ACQ_NOTE, {"size": 10.5, "italic": True, "color": GREY})]])

# =========================================================
# 16 — bo_generic_ei  (inner terms)
# =========================================================
s = slide("Active BO · EI", "bo_generic_ei — expected gain over the best-so-far")
mono(s, LX, Inches(1.5), Inches(7.5), Inches(1.5), [
    "f⁺ = max observed log-q       z = (μ − f⁺)/σ",
    "EI = (μ − f⁺)·Φ(z) + σ·φ(z)     Φ=CDF  φ=PDF",
], size=13.5, title="acquisition  (expected improvement)")
acq_table(s, LX, Inches(3.15), Inches(7.5),
          [("σ", "sigma", f2), ("z", "z_ei", f2), ("EI", "ei", lambda v: f"{v:.3f}")], winner=0, size=12)
textbox(s, LX, Inches(5.25), Inches(7.5), Inches(1.6), [
    [("(μ−f⁺)·Φ(z)", {"font": FM, "size": 13}), (" = gain × its probability;   ", {"size": 13}),
     ("σ·φ(z)", {"font": FM, "size": 13}), (" = frontier-uncertainty bonus.", {"size": 13})],
    [("here f⁺ = ", {"size": 13}), (f"{T.ACQ_BEST:.2f}", {"font": FM, "bold": True, "size": 13}),
     (";  A (high μ, modest σ) gives the largest EI.   ", {"size": 13}),
     (ACQ_NOTE, {"size": 10.5, "italic": True, "color": GREY})],
], space_after=6)
winbar(s, Inches(8.35), Inches(1.5), Inches(4.45), "picks", "A  (EI = 0.392)")
trace_table(s, Inches(8.35), Inches(2.1), Inches(4.45), T.RES["bo_ei"])
chip(s, Inches(8.35), Inches(4.0), Inches(4.45),
     [("EI finds the cliff at pick 1 (pair 26).", {"bold": True})], ACC2, size=12, h=0.5)
textbox(s, Inches(8.35), Inches(4.65), Inches(4.45), Inches(1.7), [
    [("same posterior, 3 verdicts:", {"size": 12.5, "bold": True})],
    [("σ → E", {"font": FM, "color": ACCENT}), ("   explore", {"color": GREY})],
    [("UCB → B", {"font": FM, "color": ACCENT}), ("   balance", {"color": GREY})],
    [("EI → A", {"font": FM, "color": ACCENT}), ("   exploit", {"color": GREY})],
], size=13, space_after=5)

# =========================================================
# 17 — OtterTune overview
# =========================================================
s = slide("Active · OtterTune", "ottertune — fixed-kernel GP + gradient search (GPRGD)")
mono(s, LX, Inches(1.5), Inches(7.45), Inches(2.35), [
    "K(X,X') = mag·exp(−‖ΔX‖/ℓ) + ridge·I     (fixed HPs)",
    "          ℓ=2,  mag=1,  ridge=1   —  NO MLE",
    "",
    "loss  L(x) = mu_mult·μ(x) − β·σ(x)        β grows with t",
    "β = UCB_SCALE·√(2·ln(d·t²·π²/6))",
], size=13, title="exact OtterTune kernel & acquisition  (cmu-db/ottertune)")
bullets(s, LX, Inches(4.0), Inches(7.45), Inches(2.7), [
    (0, [("Seed = ", {}), ("LHS bootstrap", {"bold": True}), (" pairs; standardize X & log q each round.", {})]),
    (0, [("Minimize L over ", {}), ("continuous", {"bold": True, "color": ACC2}),
         (" x by Adam (GPRGD), from many random + best-seen starts.", {})]),
    (0, [("Snap", {"bold": True}), (" the optimum to the nearest ", {}),
         ("un-sampled", {"bold": True}), (" pair → sample it.", {})]),
])
bullets(s, Inches(8.3), Inches(4.0), Inches(4.5), Inches(2.7), [
    (1, [("mu_mult = 1", {"font": FM, "bold": True}), (" → hunt high q  (", {}),
         ("ottertune_max", {"font": FM}), (")", {})]),
    (1, [("mu_mult = 0", {"font": FM, "bold": True}), (" → pure σ  (", {}),
         ("ottertune_reconstruct", {"font": FM}), (")", {})]),
    (2, "continuous search ≠ scoring each pair; that's the key difference from bo_generic."),
])
chip(s, Inches(8.3), Inches(1.55), Inches(4.5),
     [("β at iter 1 (toy) = ", {"bold": True}), (f"{T.RES['ot_max']['snap']['beta']:.3f}", {"font": FM})],
     ACCENT, size=12.5, h=0.5)
table(s, Inches(8.3), Inches(2.2), Inches(4.5),
      ["uncertain pair", "σ(otter)"],
      [[[(T.pair_label(pid), {"font": FM, "size": 11})], f"{T.RES['ot_max']['snap']['sig'][pid]:.2f}"]
       for pid in [24, 31, 13, 0]],
      col_w=[2.6, 1.0], size=11.5, row_h=0.32)

# =========================================================
# 18 — ottertune_max
# =========================================================
s = slide("OtterTune · max", "ottertune_max — chase the single worst pair")
mono(s, LX, Inches(1.5), Inches(7.45), Inches(1.1), [
    "minimize  L(x) = (1)·μ(x) − β·σ(x)   on  −log q",
    "  ⇒ search is pulled to high-q AND uncertain regions",
], size=13.5, title="mu_mult = 1")
bullets(s, LX, Inches(2.85), Inches(7.45), Inches(2.6), [
    (0, "Goal: maximise max_qerr_found — discover the biggest cliff."),
    (0, [("Each round it descends toward the highest predicted q, then ", {}),
         ("samples the nearest real pair.", {})]),
    (0, [("On the toy it locks onto the cliff: pick 3 = pair 24 (", {}),
         ("q = 7.61", {"bold": True, "color": ACC2}), (").", {})]),
])
winbar(s, Inches(8.35), Inches(1.5), Inches(4.45), "best toy find", "pair 24 · q = 7.61")
trace_table(s, Inches(8.35), Inches(2.1), Inches(4.45), T.RES["ot_max"])
chip(s, Inches(8.35), Inches(4.05), Inches(4.45),
     [("Strong on S_max / coverage, weaker on average accuracy.", {"bold": True})],
     ACC2, size=12, h=0.7)

# =========================================================
# 19 — ottertune_reconstruct
# =========================================================
s = slide("OtterTune · reconstruct", "ottertune_reconstruct — pure uncertainty sampling")
mono(s, LX, Inches(1.5), Inches(7.45), Inches(1.1), [
    "minimize  L(x) = (0)·μ(x) − β·σ(x)  =  −β·σ(x)",
    "  ⇒ ignore μ, just maximise σ  (cover the space)",
], size=13.5, title="mu_mult = 0")
bullets(s, LX, Inches(2.85), Inches(7.45), Inches(2.6), [
    (0, "Goal: reconstruct the whole surface, not just its peak."),
    (0, [("Behaves like ", {}), ("bo_generic_sigma", {"font": FM}),
         (", but with OtterTune's fixed kernel + GPRGD.", {})]),
    (0, "Spreads samples for coverage → better average fit, may under-find the max."),
])
winbar(s, Inches(8.35), Inches(1.5), Inches(4.45), "strategy", "spread, not peak-seek")
trace_table(s, Inches(8.35), Inches(2.1), Inches(4.45), T.RES["ot_rec"])
chip(s, Inches(8.35), Inches(4.05), Inches(4.45),
     [("max vs reconstruct = same kernel, opposite goal (μ_mult).", {"bold": True})],
     ACCENT, size=12, h=0.7)

# =========================================================
# 20 — smoothness objectives
# =========================================================
s = slide("Active · smoothness", "smoothness_* — one acquisition per target statistic", accent=ACC2)
table(s, LX, Inches(1.5), Inches(12.23),
      ["estimate", "definition", "acquisition", "intuition"],
      [
       [[("S_max", {"font": FM, "bold": True})], [("max q", {"font": FM})],
        [("GP-UCB, β_t ↑ with t", {"font": FM, "size": 11})], "concentrate on the peak"],
       [[("S_avg", {"font": FM, "bold": True})], [("mean q", {"font": FM})],
        [("Bayesian Quadrature", {"font": FM, "size": 11})], "represent the whole mean"],
       [[("S_topk", {"font": FM, "bold": True})], [("mean of top-10% q", {"font": FM})],
        [("EI over threshold τ_t", {"font": FM, "size": 11})], "concentrate on the tail"],
      ], col_w=[1.2, 2.3, 2.8, 3.4], size=12.5, header_fill=ACC2, row_h=0.55)
bullets(s, LX, Inches(3.7), Inches(12.23), Inches(2.8), [
    (0, [("All three seed with 6 LHS pairs and a Matérn GP — only the ", {}),
         ("acquisition", {"bold": True}), (" differs (like bo_generic).", {})]),
    (0, [("smoothness_max", {"font": FM, "color": ACC2}),
         (" vs ", {}), ("bo_generic_ucb", {"font": FM}),
         (": β grows with t (theory-driven) instead of a fixed κ — explores early, exploits late.", {})]),
    (0, [("smoothness_avg", {"font": FM, "color": ACC2}),
         (" vs ", {}), ("σ", {"font": FM}),
         (": picks the most ", {}), ("representative", {"italic": True}),
         (" pair (lowest mean-variance), not the most uncertain one.", {})]),
    (0, [("smoothness_topk", {"font": FM, "color": ACC2}),
         (" vs ", {}), ("EI", {"font": FM}),
         (": threshold τ_t = k-th largest observed (the tail), not the single best f⁺.", {})]),
])

# =========================================================
# 21 — smoothness_max
# =========================================================
s = slide("Smoothness · max", "smoothness_max — GP-UCB with a growing β", accent=ACC2)
mono(s, LX, Inches(1.5), Inches(7.5), Inches(1.2), [
    "acq(X) = μ(X) + β_t·σ(X)",
    "β_t = √( 2·ln( P·t²·π²/6 ) )      (Srinivas 2010)",
], size=14, title="acquisition  ( adaptive β )")
acq_table(s, LX, Inches(2.85), Inches(7.5),
          [("σ", "sigma", f2), ("μ+β·σ", "bucb", f2)], winner=1, size=12.5)
textbox(s, LX, Inches(4.95), Inches(7.5), Inches(0.5),
        [[(f"β_t ≈ {T.ACQ_BETA:.1f} at this iteration  (vs fixed κ=2 in UCB)  →  σ weighs more, exploration stays alive.  ",
           {"size": 12}), (ACQ_NOTE, {"size": 10.5, "italic": True, "color": GREY})]])
winbar(s, Inches(8.35), Inches(1.5), Inches(4.45), "picks", "B  (μ+β·σ = 4.60)", color=ACC2)
trace_table(s, Inches(8.35), Inches(2.1), Inches(4.45), T.RES["sm_max"])
chip(s, Inches(8.35), Inches(4.05), Inches(4.45),
     [("Finds cliff pairs 24 & 22 by picks 3–4.", {"bold": True})], ACC2, size=12, h=0.55)

# =========================================================
# 22 — smoothness_avg  (Bayesian Quadrature)
# =========================================================
s = slide("Smoothness · avg", "smoothness_avg — Bayesian Quadrature for the mean", accent=ACC2)
mono(s, LX, Inches(1.5), Inches(7.5), Inches(1.6), [
    "acq(p) = ( Σ_q k*(q,p) )² / σ*²(p)",
    "  k*(q,p) = posterior cross-cov of pair q & candidate p",
    "  σ*²(p)  = posterior variance at p",
], size=13, title="acquisition  ( minimise variance of the mean estimate )")
_bq = T.RES["sm_avg"]["snap"]; _cand = _bq["candidates"]
_ord = list(np.argsort(-_bq["scores"])[:5])
table(s, LX, Inches(3.25), Inches(7.5),
      ["candidate", "Σ k*(q,p)", "σ*²(p)", "score"],
      [[[(T.pair_label(_cand[k]), {"font": FM, "size": 11})],
        f"{_bq['sum_post_cov'][k]:.2f}", f"{_bq['sigma2_star'][k]:.3f}",
        [(f"{_bq['scores'][k]:.1f}", {"bold": True})]] for k in _ord],
      col_w=[2.6, 1.2, 1.1, 1.0], size=12, header_fill=ACC2, row_h=0.33,
      hi={0})
winbar(s, Inches(8.35), Inches(1.5), Inches(4.45), "picks", "central pair 14", color=ACC2)
bullets(s, Inches(8.35), Inches(3.05), Inches(4.45), Inches(1.4), [
    (1, [("high Σk*", {"font": FM}), (": correlated with many pairs (central).", {})]),
    (1, [("÷σ*²", {"font": FM}), (": discount what we're already unsure about.", {})]),
    (1, "picks the most representative pair, not the most extreme."),
])
trace_table(s, Inches(8.35), Inches(4.95), Inches(4.45), T.RES["sm_avg"])
textbox(s, LX, Inches(6.5), Inches(7.5), Inches(0.4),
        [[("real toy computation (non-degenerate: OtterTune-style kernel sums).", {"size": 10.5, "italic": True, "color": GREY})]])

# =========================================================
# 23 — smoothness_topk
# =========================================================
s = slide("Smoothness · topk", "smoothness_topk — EI over an adaptive tail threshold", accent=ACC2)
mono(s, LX, Inches(1.5), Inches(7.5), Inches(1.55), [
    "τ_t = k-th largest observed log-q   (k = 10% of P)",
    "z = (μ − τ_t)/σ",
    "acq(X) = (μ − τ_t)·Φ(z) + σ·φ(z)",
], size=13.5, title="acquisition  ( EI aimed at the top-k tail )")
acq_table(s, LX, Inches(3.2), Inches(7.5),
          [("σ", "sigma", f2), ("z", "z_tk", f2), ("EI_τ", "eitk", lambda v: f"{v:.3f}")], winner=0, size=12)
textbox(s, LX, Inches(5.3), Inches(7.5), Inches(0.5),
        [[(f"τ = {T.ACQ_TAU:.2f} (the tail floor).  Like EI, but the bar is the top-k tail, not the single best.  ",
           {"size": 12}), (ACQ_NOTE, {"size": 10.5, "italic": True, "color": GREY})]])
winbar(s, Inches(8.35), Inches(1.5), Inches(4.45), "picks", "A  (EI_τ = 0.279)", color=ACC2)
trace_table(s, Inches(8.35), Inches(2.1), Inches(4.45), T.RES["sm_topk"])
chip(s, Inches(8.35), Inches(4.05), Inches(4.45),
     [("Aims at the whole tail → best for estimating S_topk.", {"bold": True})], ACC2, size=12, h=0.7)

# =========================================================
# 24-25 — Q5 SUMMARY  (real data)
# =========================================================
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
def _load(cfg):
    p = os.path.join(ROOT, f"gt_results_sf1_10x10_{cfg}/qt5/m0/interpolation_results/summary.csv")
    d = pd.read_csv(p)
    d = d[d.method != "ground_truth"].sort_values("QERR_MEDIAN_UNSEEN").reset_index(drop=True)
    return d

def summary_slide(cfg, title, accent, takeaways):
    d = _load(cfg)
    s = slide(f"Q5 results · {cfg}", title, accent=accent)
    # winners by objective
    w_acc  = d.loc[d.QERR_MEDIAN_UNSEEN.idxmin(), "method"]
    w_r2   = d.loc[d.R2_UNSEEN.idxmax(), "method"]
    w_cov  = d.loc[d.coverage_5x_pct.idxmax(), "method"]
    w_savg = d.loc[d.S_avg_relerr.idxmin(), "method"]
    w_stk  = d.loc[d.S_topk_relerr.idxmin(), "method"]
    rows = []
    for i, r in d.iterrows():
        rows.append([str(i+1), [(r["method"], {"font": FM, "size": 10})],
                     f"{r.QERR_MEDIAN_UNSEEN:.2f}", f"{r.QERR_P90_UNSEEN:.2f}",
                     f"{r.WITHIN_5X_UNSEEN:.0f}", f"{r.R2_UNSEEN:.2f}",
                     f"{r.max_qerr_found:.1f}", f"{r.coverage_5x_pct:.0f}",
                     f"{r.S_avg_relerr:.0f}", f"{r.S_topk_relerr:.0f}"])
    table(s, LX, Inches(1.45), Inches(8.0),
          ["#", "method", "Q̃ᵤ", "P90ᵤ", "W5X", "R²ᵤ", "maxF", "c5x", "Savg%", "Stk%"],
          rows, col_w=[0.35, 2.3, 0.7, 0.7, 0.6, 0.7, 0.7, 0.6, 0.8, 0.8],
          size=10, header_fill=accent, row_h=0.305, hi={0})
    textbox(s, LX, Inches(5.5), Inches(8.0), Inches(1.4), [
        [("Q̃ᵤ", {"font": FM, "bold": True}), ("=median pred-q (unseen, ↓)   ", {"size": 10.5}),
         ("P90ᵤ", {"font": FM, "bold": True}), ("=90th pct (↓)   ", {"size": 10.5}),
         ("W5X", {"font": FM, "bold": True}), ("=%within 5× (↑)   ", {"size": 10.5}),
         ("R²ᵤ", {"font": FM, "bold": True}), ("(↑)", {"size": 10.5})],
        [("maxF", {"font": FM, "bold": True}), ("=max q found   ", {"size": 10.5}),
         ("c5x", {"font": FM, "bold": True}), ("=%of q>5 sampled (↑)   ", {"size": 10.5}),
         ("Savg%/Stk%", {"font": FM, "bold": True}), ("=rel-err of S_avg / S_topk (↓)", {"size": 10.5})],
        [(f"grid m0 (data-space uniform) · P=220 pairs · budget=22 (10%) · max true q = {d.max_qerr_true.iloc[0]:.1f}×",
          {"size": 10.5, "italic": True, "color": GREY})],
    ], space_after=4)
    # winners panel
    wins = [("accuracy  Q̃↓", w_acc), ("fit  R²↑", w_r2), ("worst-case  cov5x↑", w_cov),
            ("S_avg  err↓", w_savg), ("S_topk  err↓", w_stk)]
    textbox(s, Inches(8.65), Inches(1.45), Inches(4.15), Inches(0.35),
            [[("BEST METHOD PER GOAL", {"size": 12, "bold": True, "color": accent})]])
    yy = Inches(1.85)
    for lbl, m in wins:
        chip(s, Inches(8.65), yy, Inches(4.15),
             [(lbl+":  ", {"bold": True}), (m, {"font": FM})], accent, size=11, h=0.42)
        yy = yy + Inches(0.5)
    bullets(s, Inches(8.65), Inches(4.5), Inches(4.2), Inches(2.4), takeaways, size=12)
    return s

summary_slide("s1q2", "Parallelism ON (s1q2): who wins where", ACCENT, [
    (0, [("bo_generic_sigma", {"font": FM, "color": ACCENT}), (": best all-rounder (R²=0.60, top-3 Q̃).", {})]),
    (0, [("smoothness_topk / bo_ucb / bo_ei", {"font": FM}), (": own the worst-case (cov5x ≥ 90%).", {})]),
    (0, [("smoothness_avg", {"font": FM}), (": lowest median error here.", {})]),
    (0, [("Interpolators are fine on Q̃ but blind to cliffs (cov5x→0).", {"color": RED})]),
])

summary_slide("s1q0", "Parallelism OFF (s1q0): rougher surface", RED, [
    (0, [("Surface is rougher: max q ", {}), ("43×", {"bold": True, "color": RED}),
         (" vs 25× — tails matter more.", {})]),
    (0, [("bo_generic_sigma", {"font": FM, "color": ACCENT}), (": again most consistent (R²=0.62).", {})]),
    (0, [("bo_ucb / bo_ei / smoothness_*", {"font": FM}), (": find the 43× pair; interpolators miss it.", {})]),
    (0, [("uniform_stride wins median but has the worst tails (P90=5.2).", {"color": RED})]),
])

# ___NEXT___
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "interpolation_methods_explained.pptx")
prs.save(OUT)
print("saved", os.path.abspath(OUT), "slides:", _page[0])
