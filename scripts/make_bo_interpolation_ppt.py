#!/usr/bin/env python3
# =========================================================
# scripts/make_bo_interpolation_ppt.py
# ---------------------------------------------------------
# Builds  bo_interpolation_explained.pptx  — a formula-first
# deck explaining scripts/bo_interpolation.py from scratch:
# Bayesian Optimization as an online sampler that hunts the
# high-q-error (boundary) pairs of a fixed query/DB.
#
# Tables / formulas / plots all come from the real run on
# gt_results_sf1_10x10_s1q0 (bo_interpolation + evaluate).
#
# Run:  python3 scripts/make_bo_interpolation_ppt.py [gt_root]
# =========================================================

import sys
from pathlib import Path

import numpy as np
import pandas as pd

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.lines import Line2D

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------
# palette / fonts
# ---------------------------------------------------------
INK    = RGBColor(0x1F, 0x2A, 0x37)
ACCENT = RGBColor(0x0E, 0x7C, 0x86)   # teal
ACC2   = RGBColor(0xB4, 0x53, 0x09)   # amber (winner / highlight)
GREEN  = RGBColor(0x15, 0x6F, 0x3B)
RED    = RGBColor(0xB1, 0x2A, 0x2A)
GREY   = RGBColor(0x6B, 0x72, 0x80)
LIGHT  = RGBColor(0xF2, 0xF4, 0xF6)
ROW1   = RGBColor(0xFF, 0xFF, 0xFF)
ROW2   = RGBColor(0xEC, 0xF3, 0xF4)
HIROW  = RGBColor(0xFD, 0xF1, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BAND   = RGBColor(0xE8, 0xEC, 0xEF)
FT, FM = "Segoe UI", "Consolas"

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
BLANK = prs.slide_layouts[6]
_page = [0]
FOOTER = "Bayesian Optimization · q-error boundary discovery  ·  scripts/bo_interpolation.py"


# ---------------------------------------------------------
# low-level helpers  (same toolkit as make_interpolation_ppt.py)
# ---------------------------------------------------------
def _set_fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def _runs(p, runs, default_size, default_color, font=FT, align=None):
    if align is not None:
        p.alignment = align
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, o in runs:
        r = p.add_run(); r.text = text; f = r.font
        f.name = o.get("font", font); f.size = Pt(o.get("size", default_size))
        f.bold = o.get("bold", False); f.italic = o.get("italic", False)
        f.color.rgb = o.get("color", default_color)
    return p

def textbox(slide, l, t, w, h, lines, size=16, color=INK, font=FT,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=4, line_spacing=1.0):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after); p.space_before = Pt(0); p.line_spacing = line_spacing
        _runs(p, ln, size, color, font=font, align=align)
    return tb

def rrect(slide, l, t, w, h, fill, line=None, line_w=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sp.adjustments[0] = 0.06
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def rect(slide, l, t, w, h, fill):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    _set_fill(sp, fill); sp.shadow.inherit = False
    return sp

def slide(kicker, title, accent=ACCENT):
    s = prs.slides.add_slide(BLANK); _page[0] += 1
    rect(s, 0, 0, SW, Inches(0.16), accent)
    textbox(s, Inches(0.55), Inches(0.24), Inches(11.5), Inches(0.3),
            [[(kicker.upper(), {"size": 12.5, "bold": True, "color": accent})]])
    textbox(s, Inches(0.55), Inches(0.52), Inches(12.3), Inches(0.8),
            [[(title, {"size": 26, "bold": True, "color": INK})]])
    rect(s, Inches(0.55), Inches(1.26), Inches(12.23), Pt(1.6), BAND)
    textbox(s, Inches(0.55), Inches(7.12), Inches(10.5), Inches(0.3),
            [[(FOOTER, {"size": 9, "color": GREY})]])
    textbox(s, Inches(12.2), Inches(7.12), Inches(0.9), Inches(0.3),
            [[(str(_page[0]), {"size": 9, "color": GREY})]], align=PP_ALIGN.RIGHT)
    return s

def mono(slide, l, t, w, h, lines, size=14, title=None, fill=LIGHT,
         color=INK, accent=ACCENT, line_spacing=1.05):
    rrect(slide, l, t, w, h, fill, line=BAND, line_w=1.0)
    pad = Inches(0.14); ty = t + Inches(0.06)
    if title:
        textbox(slide, l + pad, ty, w - 2*pad, Inches(0.3),
                [[(title, {"size": size-1.5, "bold": True, "color": accent, "font": FT})]])
        ty = ty + Inches(0.34)
    textbox(slide, l + pad, ty, w - 2*pad, h - (ty - t) - Inches(0.06),
            lines, size=size, color=color, font=FM, line_spacing=line_spacing, space_after=3)

def bullets(slide, l, t, w, h, items, size=15.5, color=INK, accent=ACCENT, gap=5, line_spacing=1.02):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    marks = {0: "▸  ", 1: "–  ", 2: "·  "}
    for i, it in enumerate(items):
        lvl, runs = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(gap); p.space_before = Pt(0); p.line_spacing = line_spacing
        mcol = accent if lvl == 0 else (GREY if lvl == 2 else ACCENT)
        runs2 = [(marks.get(lvl, ""), {"color": mcol, "bold": True})]
        if isinstance(runs, str):
            runs2.append((runs, {}))
        else:
            runs2.extend(runs)
        _runs(p, runs2, size if lvl == 0 else size-1.5, color)
    return tb

def _cell(cell, runs, size, color, bold, align, fill, font=FT):
    cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    if isinstance(runs, str): runs = [(runs, {})]
    for text, o in runs:
        r = p.add_run(); r.text = text; f = r.font
        f.name = o.get("font", font); f.size = Pt(o.get("size", size))
        f.bold = o.get("bold", bold); f.italic = o.get("italic", False)
        f.color.rgb = o.get("color", color)

def table(slide, l, t, w, header, rows, col_w=None, size=12, header_fill=ACCENT,
          header_color=WHITE, hi=None, hi_fill=HIROW, aligns=None, font=FT,
          row_h=0.34, header_h=0.36, body_color=INK):
    nC = len(header); nR = len(rows) + 1
    h = Inches(header_h + row_h * len(rows))
    gt = slide.shapes.add_table(nR, nC, l, t, w, h).table
    gt.first_row = False; gt.horz_banding = False
    tblPr = gt._tbl.tblPr; tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')
    if col_w:
        tot = sum(col_w)
        for j, cw in enumerate(col_w):
            gt.columns[j].width = Emu(int(w * cw / tot))
    gt.rows[0].height = Inches(header_h)
    if aligns is None:
        aligns = [PP_ALIGN.LEFT] + [PP_ALIGN.CENTER]*(nC-1)
    for j, htxt in enumerate(header):
        _cell(gt.cell(0, j), htxt, size, header_color, True, aligns[j], header_fill, font=font)
    hi = hi or set()
    for i, row in enumerate(rows):
        gt.rows[i+1].height = Inches(row_h)
        base = hi_fill if i in hi else (ROW2 if i % 2 else ROW1)
        for j, c in enumerate(row):
            _cell(gt.cell(i+1, j), c, size, body_color, False, aligns[j], base, font=font)
    return gt

def chip(slide, l, t, w, text, fill, tcolor=WHITE, size=12.5, h=0.38):
    sp = rrect(slide, l, t, w, Inches(h), fill)
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _runs(p, text if isinstance(text, list) else [(text, {})], size, tcolor, font=FT)
    return sp

def pic(slide, path, l, t, w=None, h=None):
    slide.shapes.add_picture(str(path), l, t, width=w, height=h)

LX = Inches(0.55); CW = Inches(12.23)
def f2(v): return f"{v:.2f}"
def f3(v): return f"{v:.3f}"


# =========================================================
# DATA  (real run)
# =========================================================
GT_ROOT   = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("gt_results_sf1_10x10_s1q0")
EVAL_CSV  = GT_ROOT / "bo_interpolation_eval" / "bo_eval_all.csv"
DEMO_DIR  = GT_ROOT / "qt8" / "m0" / "bo_interpolation_results"
ASSET_DIR = Path("bo_ppt_assets"); ASSET_DIR.mkdir(exist_ok=True)
OUT_PPTX  = Path("bo_interpolation_explained.pptx")

df = pd.read_csv(EVAL_CSV)
P  = "QERR_MEDIAN_UNSEEN"
noref = df[df.method != "ground_truth"].copy()
bo    = noref[noref.method.str.startswith("bo_")].copy()

n_q   = df["query"].nunique()
n_md  = df.groupby(["query", "gt_method"]).ngroups
n_meth= df["method"].nunique()

agg = (noref.groupby(["method", "representation", "kernel", "acquisition"])
            .agg(qmed=(P, "mean"), p90=("QERR_P90_UNSEEN", "mean"),
                 w2=("WITHIN_2X_UNSEEN", "mean"), maxf=("max_found_pct", "mean"),
                 cov5=("coverage_5x_pct", "mean"))
            .reset_index().sort_values("qmed").reset_index(drop=True))

def _grp(col):
    return (bo.groupby(col).agg(qmed=(P, "mean"), maxf=("max_found_pct", "mean"),
                                cov5=("coverage_5x_pct", "mean"))
              .reset_index().sort_values("qmed").reset_index(drop=True))

by_acq  = _grp("acquisition")
by_kern = _grp("kernel")
by_rep  = _grp("representation")

best_overall = agg.iloc[0]
best_bo      = agg[agg.method.str.startswith("bo_")].iloc[0]
best_disc    = agg.sort_values("maxf", ascending=False).iloc[0]
budget_pct   = 100.0 * noref["sample_fraction"].mean()


# =========================================================
# PLOTS
# =========================================================
plt.rcParams.update({"font.size": 11, "axes.grid": True, "grid.alpha": 0.25,
                     "figure.dpi": 150})

def plot_tradeoff(path):
    fig, ax = plt.subplots(figsize=(6.6, 4.6))
    acqs = ["sigma", "ts", "ucb", "pi", "ei"]
    cmap = {"sigma": "#0E7C86", "ts": "#7B68EE", "ucb": "#B45309",
            "pi": "#B12A2A", "ei": "#156F3B"}
    for a in acqs:
        sub = bo[bo.acquisition == a]
        ax.scatter(sub.coverage_5x_pct, sub.max_found_pct, s=70, alpha=0.8,
                   color=cmap[a], edgecolors="white", linewidths=0.6, label=a)
    ax.set_xlabel("coverage of q>5x pairs  (%)   → breadth")
    ax.set_ylabel("max q-error found  (% of true max)   → peak")
    ax.set_title("Boundary discovery: every BO config\n(colored by acquisition)")
    ax.legend(title="acquisition", fontsize=9, loc="lower right")
    fig.tight_layout(); fig.savefig(path); plt.close(fig)

def plot_by_acq(path):
    fig, axs = plt.subplots(1, 2, figsize=(9.2, 4.2))
    a = by_acq.sort_values("qmed")
    axs[0].bar(a.acquisition, a.qmed, color="#0E7C86")
    axs[0].axhline(df[df.method == "random"][P].mean(), ls="--", color="#B12A2A", lw=1.5,
                   label="random baseline")
    axs[0].set_title("Interpolation accuracy\n(median unseen q-error, lower better)")
    axs[0].legend(fontsize=8)
    a2 = by_acq.sort_values("cov5", ascending=False)
    x = np.arange(len(a2)); wdt = 0.4
    axs[1].bar(x - wdt/2, a2.maxf, wdt, color="#B45309", label="max found %")
    axs[1].bar(x + wdt/2, a2.cov5, wdt, color="#156F3B", label="coverage>5x %")
    axs[1].set_xticks(x); axs[1].set_xticklabels(a2.acquisition)
    axs[1].set_title("Boundary discovery\n(higher better)"); axs[1].legend(fontsize=8)
    fig.tight_layout(); fig.savefig(path); plt.close(fig)

def plot_by_kern_rep(path):
    fig, axs = plt.subplots(1, 2, figsize=(9.2, 4.0))
    axs[0].bar(by_kern.kernel, by_kern.qmed, color="#0E7C86")
    axs[0].set_title("By kernel — median unseen q-error"); axs[0].tick_params(axis="x", rotation=20)
    axs[1].bar(by_rep.representation, by_rep.qmed, color="#B45309", width=0.5)
    axs[1].set_title("By representation — median unseen q-error")
    for ax in axs:
        ax.set_ylabel("median unseen q-error")
    fig.tight_layout(); fig.savefig(path); plt.close(fig)

def plot_surface(path):
    """qt8/m0: where sigma (explore) vs ucb (exploit) place their samples."""
    gt = pd.read_csv(DEMO_DIR / "ground_truth" / "predictions.csv")
    mx1 = (gt.a_x1 + gt.b_x1) / 2.0; mx2 = (gt.a_x2 + gt.b_x2) / 2.0
    logq = np.log10(np.maximum(gt.y_true.values, 1.0))
    fig, axs = plt.subplots(1, 2, figsize=(10.4, 4.7))
    for ax, meth, ttl in [(axs[0], "bo_pair_matern52_sigma", "sigma  (explore)"),
                          (axs[1], "bo_pair_matern52_ucb",   "ucb  (exploit)")]:
        pv = pd.read_csv(DEMO_DIR / meth / "predictions.csv")
        sc = ax.scatter(mx1, mx2, c=logq, cmap="viridis", s=26)
        smp = pv[pv.is_sampled == 1]
        smx1 = (smp.a_x1 + smp.b_x1) / 2.0; smx2 = (smp.a_x2 + smp.b_x2) / 2.0
        ax.scatter(smx1, smx2, s=150, facecolors="none", edgecolors="red", linewidths=1.8,
                   label="sampled")
        mf = 100.0 * smp.y_true.max() / gt.y_true.max()
        n5 = (gt.y_true > 5).sum(); c5 = 100.0 * (smp.y_true > 5).sum() / max(n5, 1)
        ax.set_title(f"{ttl}\nmax found {mf:.0f}%   ·   coverage>5x {c5:.0f}%")
        ax.set_xlabel("mid x1"); ax.set_ylabel("mid x2"); ax.legend(fontsize=9, loc="upper right")
    cb = fig.colorbar(sc, ax=axs, fraction=0.046, pad=0.02)
    cb.set_label("log₁₀ true q-error")
    fig.suptitle("Same 10×10 grid (qt8/m0, matern5/2): sampling pattern by acquisition", y=1.02)
    fig.savefig(path, bbox_inches="tight"); plt.close(fig)

P_TRADE = ASSET_DIR / "tradeoff.png"
P_ACQ   = ASSET_DIR / "by_acq.png"
P_KR    = ASSET_DIR / "by_kern_rep.png"
P_SURF  = ASSET_DIR / "surface.png"
plot_tradeoff(P_TRADE); plot_by_acq(P_ACQ); plot_by_kern_rep(P_KR); plot_surface(P_SURF)
print("plots written to", ASSET_DIR)


# =========================================================
# 1 — TITLE
# =========================================================
s = prs.slides.add_slide(BLANK); _page[0] += 1
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, Inches(4.55), SW, Inches(0.05), ACCENT)
textbox(s, Inches(0.8), Inches(1.4), Inches(11.9), Inches(2.6), [
    [("Bayesian Optimization for ", {"size": 40, "bold": True, "color": WHITE}),
     ("Q-error Boundaries", {"size": 40, "bold": True, "color": RGBColor(0x6F, 0xD3, 0xDA)})],
    [("Using BO as an online sampler to find where the optimizer mis-predicts — not to predict runtime",
      {"size": 18, "color": RGBColor(0xC8, 0xD2, 0xDA)})],
], space_after=10)
textbox(s, Inches(0.8), Inches(4.8), Inches(11.9), Inches(1.7), [
    [("Objective:   ", {"size": 17, "color": RGBColor(0x9F, 0xB0, 0xBC)}),
     ("maximise q(P₁,P₂) = max(rt₁,rt₂)/min(rt₁,rt₂)  over predicate values, fixed template + DB",
      {"size": 17, "bold": True, "color": WHITE, "font": FM})],
    [(f"3 representations · 5 kernels · 5 acquisitions (+3 baselines) = {n_meth} methods · "
      f"10% budget · run on {GT_ROOT.name}  ({n_q} queries × 3 sampling methods = {n_md} grids)",
      {"size": 13.5, "color": RGBColor(0x9F, 0xB0, 0xBC)})],
], space_after=12)
textbox(s, Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
        [[("scripts/bo_interpolation.py  +  evaluate_bo_interpolation.py  +  bo_interpolation_summary.py",
           {"size": 12, "color": ACCENT, "font": FM})]])

# =========================================================
# 2 — OBJECTIVE / FORMULATION
# =========================================================
s = slide("Problem framing", "The goal is a sampler, not a runtime model")
bullets(s, LX, Inches(1.5), Inches(6.15), Inches(4.4), [
    (0, [("We do ", {}), ("not", {"bold": True, "color": RED}), (" build a general runtime predictor.", {})]),
    (0, [("We use BO as an ", {}), ("online intelligent sampler", {"bold": True}),
         (" to efficiently discover ", {}), ("high-q-error (boundary) regions", {"bold": True, "color": ACC2}), (".", {})]),
    (0, [("Quantity of interest is ", {}), ("q-error", {"bold": True}), (", not runtime itself.", {})]),
    (0, "Condition on everything except the predicate values:"),
    (2, "query template, DB instance, indexes, PG config, hardware — all fixed."),
    (0, [("So the black box reduces to  ", {}),
         ("q = f(predicate values) = f(grid point)", {"font": FM, "bold": True, "color": ACCENT}), (".", {})]),
    (0, [("Another DB / template = a ", {}), ("new", {"italic": True}),
         (" optimization problem (no cross-DB transfer).", {})]),
])
mono(s, Inches(6.95), Inches(1.5), Inches(5.85), Inches(2.55), [
    "initial random samples",
    "      │",
    "      ▼   measure q  →  fit GP surrogate",
    "      ▼   acquisition picks next point",
    "      ▼   measure q  →  update GP",
    "      └────────── repeat (online) ─────────┘",
], size=13, title="online BO — never build the full grid first")
mono(s, Inches(6.95), Inches(4.25), Inches(5.85), Inches(1.15), [
    "q = f( predicate_values | template, DB,",
    "        indexes, pg_config, hardware )",
], size=13, title="the conditioned black box", accent=ACC2)
textbox(s, Inches(6.95), Inches(5.55), Inches(5.85), Inches(1.2), [
    [("Analogous to BO for hyper-parameter tuning: dataset & model fixed, only the knobs "
      "(here, predicate values) vary.", {"size": 12, "italic": True, "color": GREY})]])

# =========================================================
# 3 — LABEL + PAIRS
# =========================================================
s = slide("Foundations · label & instances", "Instance = axis-neighbour pair;  label = its q-error")
mono(s, LX, Inches(1.5), Inches(6.0), Inches(1.2), [
    "q(a,b) = max(rtₐ, rt_b) / min(rtₐ, rt_b)  ≥ 1",
    "model in log-space:  z = log q   →   q̂ = eᶻ ≥ 1",
], size=14, title="symmetric ratio  (bo_interpolation.symmetric_ratio)")
mono(s, LX, Inches(2.95), Inches(6.0), Inches(1.75), [
    "points        = N × N        = 11 × 11 = 121",
    "pairs / axis  = (N−1) × N    = 10 × 11 = 110",
    "TOTAL pairs P = 2·(N−1)·N    = 220",
    "neighbour of (i,j):  (i+1,j) [axis0] , (i,j+1) [axis1]",
], size=13, title="grid → directed axis-neighbour pairs  (10×10 resolution)")
bullets(s, Inches(6.95), Inches(1.5), Inches(5.85), Inches(5.0), [
    (0, [("q ≈ 1", {"bold": True, "color": GREEN}), (" = neighbours run alike (smooth).", {})]),
    (0, [("q large", {"bold": True, "color": ACC2}), (" = a runtime cliff between neighbours (e.g. a plan flip) — the ", {}),
         ("boundary", {"bold": True}), (" we want to find.", {})]),
    (0, [("The neighbour is ", {}), ("deterministic", {"bold": True}),
         (" once the grid is built — this enables representation B.", {})]),
    (0, [("Most of the 220 pairs are ≈1; only a handful are large.  On this grid q ranges ", {}),
         ("1 → ~1500", {"bold": True, "font": FM, "color": ACC2}), (".", {})]),
    (0, [("Budget = ", {}), ("10% of P", {"bold": True}), ("  ⇒  measure only ~22 of 220 pairs.", {})]),
])

# =========================================================
# 4 — TWO REPRESENTATIONS
# =========================================================
s = slide("Design · input representation", "Two representations of a BO input  (A is the default)")
table(s, LX, Inches(1.45), Inches(12.23),
      ["", "input", "feature X  (2-D grid)", "fdim", "pros / cons"],
      [
       [[("A · pair", {"bold": True, "color": ACC2})], [("(P₁, P₂)", {"font": FM})],
        [("[P₁ , P₂]", {"font": FM, "size": 11})], "4",
        [("general, no neighbour assumption · higher-dim, redundant", {"size": 11})]],
       [[("B · point", {"bold": True, "color": ACCENT})], [("P  (+ det. neighbour)", {"font": FM})],
        [("[P₁]", {"font": FM, "size": 11})], "2",
        [("lower-dim · H+V pairs share X=P (collision)", {"size": 11})]],
       [[("B-split · point/axis", {"bold": True, "color": GREEN})], [("P  (one BO per axis)", {"font": FM})],
        [("[P₁]  ×2 models", {"font": FM, "size": 11})], "2",
        [("recovers axis info, NO extra feature · no collision", {"size": 11})]],
      ], col_w=[1.6, 2.2, 3.3, 0.6, 4.6], size=12, row_h=0.5, hi={0},
      aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT])
mono(s, LX, Inches(3.55), Inches(6.0), Inches(1.15), [
    "both:  label = q(P₁, P₂)   (same 220 pairs)",
    "differ only in the feature X the GP sees",
], size=13, title="same target, different X")
bullets(s, Inches(6.95), Inches(3.55), Inches(5.85), Inches(3.0), [
    (0, [("Professor's hint:", {"bold": True}), (" if the neighbour is deterministic, a function of ", {}),
         ("P₁ alone", {"bold": True}), (" (rep B) may suffice.", {})]),
    (0, "All horizontal & vertical pairs are still evaluated — B only shrinks the feature, not the pair set."),
    (0, [("Test both; if similar, prefer the ", {}), ("lower-dimensional", {"bold": True}),
         (" rep B (easier GP fit).", {})]),
])
textbox(s, LX, Inches(4.85), Inches(6.0), Inches(1.4),
        [[("Feature X is exactly the specified input — no axis one-hot.  Rep B lets a point's H & V "
           "pairs share X = P; B-split fixes that with one BO per axis (still input = P).",
           {"size": 11.5, "italic": True, "color": GREY})]])

# =========================================================
# 5 — GP SURROGATE
# =========================================================
s = slide("Surrogate", "Gaussian-Process surrogate over log q")
mono(s, LX, Inches(1.5), Inches(6.0), Inches(2.7), [
    "fit on sampled (X, log q):",
    "  μ(X)  = k*ᵀ (K + σ²I)⁻¹ y      posterior mean",
    "  σ²(X) = k(X,X) − k*ᵀ(K+σ²I)⁻¹k*  posterior var",
    "",
    "k*  = kernel from X to every sampled pair",
    "hyper-params by MLE (sklearn L-BFGS), normalize_y",
], size=13, title="sklearn GaussianProcessRegressor  (make_gp)")
bullets(s, Inches(6.95), Inches(1.5), Inches(5.85), Inches(5.0), [
    (0, [("The GP returns both a ", {}), ("best guess μ", {"bold": True}),
         (" and an ", {}), ("uncertainty σ", {"bold": True}), (" at every un-sampled pair.", {})]),
    (0, [("σ", {"bold": True}), (" is exactly what lets an acquisition ", {}),
         ("choose", {"italic": True}), (" the next pair to measure.", {})]),
    (0, [("Everything is done in ", {}), ("log-q", {"bold": True}),
         (" space (ratios are multiplicative), then exp() back — so q̂ ≥ 1 always.", {})]),
    (0, [("kernel = ", {}), ("C·k(r) + WhiteKernel", {"font": FM}),
         (" — a signal scale, the shape below, and a learned noise floor.", {})]),
])

# =========================================================
# 6 — KERNELS
# =========================================================
s = slide("Surrogate · kernels", "Every kernel we sweep  (k(r), r = ‖X − X'‖ z-scored)")
table(s, LX, Inches(1.45), Inches(12.23),
      ["kernel", "k(r)", "smoothness / behaviour"],
      [
       [[("rbf  (SE)", {"bold": True})], [("exp(−r²/2ℓ²)", {"font": FM})],
        "∞-differentiable — very smooth (can over-smooth cliffs)"],
       [[("matern12  (ν=½)", {"bold": True})], [("exp(−r/ℓ)", {"font": FM})],
        "non-differentiable — rough, reacts fast to local jumps"],
       [[("matern32  (ν=3/2)", {"bold": True})], [("(1 + √3·r/ℓ)·exp(−√3·r/ℓ)", {"font": FM})],
        "1× differentiable"],
       [[("matern52  (ν=5/2)", {"bold": True, "color": ACC2})],
        [("(1 + √5·r/ℓ + 5r²/3ℓ²)·exp(−√5·r/ℓ)", {"font": FM})],
        [("2× differentiable — curvature but finite at a cliff (best here)", {"bold": True})]],
       [[("rquad  (RQ)", {"bold": True})], [("(1 + r²/2αℓ²)^(−α)", {"font": FM})],
        "scale-mixture of RBFs (many length-scales at once)"],
      ], col_w=[2.2, 4.3, 5.7], size=12, row_h=0.6, hi={3},
      aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT])
textbox(s, LX, Inches(6.05), Inches(12.23), Inches(0.9), [
    [("Each kernel is wrapped as  ", {"size": 12.5}),
     ("ConstantKernel(σ_f²)·k(r) + WhiteKernel(noise)", {"font": FM, "size": 12.5, "bold": True}),
     ("  and its (σ_f, ℓ, noise) fit by marginal-likelihood at every BO step.", {"size": 12.5})]])

# =========================================================
# 7 — ACQUISITIONS
# =========================================================
s = slide("Sampler · acquisitions", "Every acquisition — the one line that changes", accent=ACC2)
table(s, LX, Inches(1.4), Inches(12.23),
      ["acquisition", "score a(X)  — pick argmax", "intent"],
      [
       [[("sigma", {"bold": True})], [("σ(X)", {"font": FM})],
        [("pure explore / reconstruct", {})]],
       [[("ucb", {"bold": True, "color": ACC2})], [("μ + κ·σ      (κ = 2)", {"font": FM})],
        [("balance — hunt the peak", {})]],
       [[("ei", {"bold": True})], [("(μ−f⁺)·Φ(z) + σ·φ(z),   z=(μ−f⁺)/σ", {"font": FM, "size": 11})],
        [("expected gain over best f⁺", {})]],
       [[("pi", {"bold": True})], [("Φ(z),   z=(μ−f⁺)/σ", {"font": FM})],
        [("probability of beating f⁺", {})]],
       [[("ts  (Thompson)", {"bold": True})], [("draw g ~ GP posterior;  a = g(X)", {"font": FM})],
        [("randomised explore/exploit", {})]],
      ], col_w=[2.2, 5.6, 3.5], size=12.5, row_h=0.52, hi={1},
      aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT])
textbox(s, LX, Inches(4.7), Inches(12.23), Inches(0.5), [
    [("Φ", {"font": FM, "bold": True}), (" = normal CDF   ", {}),
     ("φ", {"font": FM, "bold": True}), (" = normal PDF   ", {}),
     ("f⁺", {"font": FM, "bold": True}), (" = best log-q seen so far   ", {}),
     ("μ,σ", {"font": FM, "bold": True}), (" = GP posterior (log-q)", {})]])
bullets(s, LX, Inches(5.25), Inches(12.23), Inches(1.6), [
    (0, [("sigma / ts", {"bold": True, "color": ACCENT}), (" spread out → good for ", {}),
         ("reconstruction", {"bold": True}), (" (accurate surface everywhere).", {})]),
    (0, [("ucb / ei / pi", {"bold": True, "color": ACC2}), (" chase high μ → good for ", {}),
         ("boundary discovery", {"bold": True}), (" (find & cover the cliffs).", {})]),
])

# =========================================================
# 8 — LOOP + BUDGET
# =========================================================
s = slide("Sampler · the online loop", "One loop, swap the acquisition; 10% budget")
mono(s, LX, Inches(1.5), Inches(7.4), Inches(2.6), [
    "seed  ← 6 random distinct pairs      (shared across all runs)",
    "while |sampled| < BUDGET:            # BUDGET = ⌈0.10·P⌉ = 22",
    "    fit GP(kernel) on sampled (X, log q)",
    "    μ, σ ← GP.predict(all 220 pairs)",
    "    score ← ACQUISITION(μ, σ)        ◄ the ONE differing line",
    "    next  ← argmax score  (un-sampled)",
    "    measure q(next) ; add to sampled  # 'measure' = read ground_truth",
], size=12.5, title="run_bo()  — refit after every observation (online)")
bullets(s, Inches(8.1), Inches(1.6), Inches(4.7), Inches(3.0), [
    (0, [("Same 6-pair seed", {"bold": True}), (" for every (rep × kernel × acq) → fair comparison.", {})]),
    (0, [("Baselines", {"bold": True}), (": random & uniform-stride sample once, then linear-interpolate.", {})]),
    (0, [("ground_truth", {"font": FM}), (" measures all 220 — the reference.", {})]),
])
mono(s, LX, Inches(4.35), Inches(7.4), Inches(1.0), [
    "MIN = max(5, 2·dim+2)=6      BUDGET = min(P, max(MIN, ⌈0.10·P⌉))",
], size=13, title="budget", accent=ACC2)
textbox(s, Inches(8.1), Inches(4.75), Inches(4.7), Inches(1.6),
        [[("Grid: ", {"bold": True}), (f"{n_meth} methods", {"bold": True, "color": ACC2}),
          (" = 3 baselines + 3 reps × 5 kernels × 5 acquisitions (75 BO), run on each of ", {}),
          (f"{n_md} grids", {"bold": True}), (".", {})]], size=13)

# =========================================================
# 9 — METRICS
# =========================================================
s = slide("Scoring", "Two families of metric: accuracy vs discovery", accent=ACC2)
mono(s, LX, Inches(1.5), Inches(6.05), Inches(1.55), [
    "pred q-error:  Q = max( y/ŷ , ŷ/y ) ≥ 1",
    "QERR_MEDIAN_UNSEEN = median Q over UN-sampled pairs",
    "WITHIN_kX = % of pairs with Q ≤ k",
], size=13, title="A. interpolation accuracy  (uses ŷ vs y on unseen)", accent=ACCENT)
mono(s, LX, Inches(3.2), Inches(6.05), Inches(2.05), [
    "max_qerr_found = max( y_true over SAMPLED pairs )",
    "max_found_pct  = 100 · max_found / max_true",
    "",
    "coverage_5x = 100 · #(sampled, y>5) / #(all, y>5)",
    "  (breadth: fraction of q>5 pairs actually sampled)",
], size=12.5, title="B. boundary discovery  (uses y_true of sampled only)", accent=ACC2)
bullets(s, Inches(6.95), Inches(1.5), Inches(5.85), Inches(5.2), [
    (0, [("Accuracy", {"bold": True, "color": ACCENT}), (" = how good is the interpolated q-surface where we ", {}),
         ("didn't", {"italic": True}), (" sample.", {})]),
    (0, [("Discovery", {"bold": True, "color": ACC2}), (" = did the sampler go to the right places?", {})]),
    (1, [("max_found", {"font": FM}), (" = did we hit the single ", {}), ("tallest", {"bold": True}), (" cliff?", {})]),
    (1, [("coverage", {"font": FM}), (" = did we catch ", {}), ("many", {"bold": True}), (" of the cliffs?", {})]),
    (0, [("These are ", {}), ("independent", {"bold": True}),
         (": a method can hit the peak yet miss most of the boundary, or vice-versa.", {})]),
    (0, "That independence is the whole story of the results →"),
])

# =========================================================
# 10 — RESULTS: ranking
# =========================================================
s = slide("Results · ranking", f"Top methods on {GT_ROOT.name}  (mean over {n_md} grids)", accent=ACC2)
rows = []
top = agg.head(9)
for _, r in top.iterrows():
    is_base = r.method in ("random", "uniform_stride")
    rows.append([
        [(r.method, {"font": FM, "size": 10.5, "bold": True,
                     "color": (GREY if is_base else INK)})],
        r.representation, r.kernel, r.acquisition,
        [(f3(r.qmed), {"bold": True})], f2(r.maxf), f2(r.cov5),
    ])
hi = {i for i, (_, r) in enumerate(top.iterrows()) if r.method.startswith("bo_")}
table(s, LX, Inches(1.5), Inches(12.23),
      ["method", "rep", "kernel", "acq", "med-unseen q ↓", "max found % ↑", "cover>5x % ↑"],
      rows, col_w=[3.4, 1.1, 1.4, 0.9, 1.9, 1.7, 1.7], size=11.5, row_h=0.42,
      header_fill=ACC2, aligns=[PP_ALIGN.LEFT]+[PP_ALIGN.CENTER]*6)
textbox(s, LX, Inches(5.9), Inches(12.23), Inches(1.0), [
    [("Baselines (grey) still win pure accuracy", {"bold": True}),
     (" — BO with the ", {}), ("sigma", {"font": FM, "bold": True}),
     (" acquisition (matern32/52) is the first to tie them.  ", {}),
     ("Ranking by median-unseen q-error rewards reconstruction, not boundary-hunting.",
      {"italic": True, "color": GREY})]])

# =========================================================
# 11 — RESULTS: by acquisition
# =========================================================
s = slide("Results · acquisition", "Acquisition drives the explore/exploit split", accent=ACC2)
pic(s, P_ACQ, LX, Inches(1.5), w=Inches(7.3))
table(s, Inches(8.1), Inches(1.6), Inches(4.7),
      ["acq", "med-q ↓", "found% ↑", "cov5% ↑"],
      [[[(r.acquisition, {"font": FM, "bold": True})], f3(r.qmed), f2(r.maxf), f2(r.cov5)]
       for _, r in by_acq.iterrows()],
      col_w=[1.3, 1.2, 1.2, 1.2], size=12, header_fill=ACC2, row_h=0.44)
bullets(s, Inches(8.1), Inches(4.4), Inches(4.7), Inches(2.4), [
    (0, [("sigma", {"font": FM, "bold": True, "color": ACCENT}), (" wins accuracy, low coverage.", {})]),
    (0, [("ucb/pi/ei", {"font": FM, "bold": True, "color": ACC2}), (" ~2× the coverage of the boundary.", {})]),
    (0, "Acquisition matters more than kernel or representation."),
])

# =========================================================
# 12 — RESULTS: kernel & representation
# =========================================================
s = slide("Results · kernel & representation", "Axis-split wins; kernel matters less", accent=ACC2)
pic(s, P_KR, LX, Inches(1.5), w=Inches(7.3))
table(s, Inches(8.1), Inches(1.6), Inches(4.7),
      ["kernel", "med-q ↓", "found% ↑"],
      [[[(r.kernel, {"font": FM, "bold": True})], f3(r.qmed), f2(r.maxf)] for _, r in by_kern.iterrows()],
      col_w=[1.6, 1.2, 1.2], size=12, header_fill=ACC2, row_h=0.4)
table(s, Inches(8.1), Inches(4.35), Inches(4.7),
      ["rep", "med-q ↓", "found% ↑", "cov5% ↑"],
      [[[(r.representation, {"font": FM, "bold": True})], f3(r.qmed), f2(r.maxf), f2(r.cov5)]
       for _, r in by_rep.iterrows()],
      col_w=[1.4, 1.2, 1.2, 1.2], size=12, header_fill=ACCENT, row_h=0.4)
textbox(s, Inches(8.1), Inches(5.7), Inches(4.7), Inches(1.3),
        [[("point_split", {"bold": True, "size": 12, "color": GREEN}),
          (" (one BO per axis) beats both pair and point on accuracy AND discovery — same input P, no collision.",
           {"size": 12})]])

# =========================================================
# 13 — VIS: sampling pattern
# =========================================================
s = slide("Results · visualization", "Where the budget lands: explore vs exploit", accent=ACC2)
pic(s, P_SURF, LX, Inches(1.45), w=Inches(9.2))
bullets(s, Inches(10.0), Inches(1.7), Inches(2.9), Inches(4.8), [
    (0, [("Same grid", {"bold": True}), (" (qt8/m0), same kernel (matern5/2), same seed.", {})]),
    (0, [("sigma", {"font": FM, "bold": True, "color": ACCENT}), (" spreads samples to reduce uncertainty everywhere.", {})]),
    (0, [("ucb", {"font": FM, "bold": True, "color": ACC2}), (" clusters on the bright high-q ridge.", {})]),
    (0, "Both hit the peak; ucb covers far more of the boundary."),
])

# =========================================================
# 14 — VIS: trade-off scatter
# =========================================================
s = slide("Results · the trade-off", "Peak-hitting vs breadth — one point per BO config", accent=ACC2)
pic(s, P_TRADE, LX, Inches(1.5), w=Inches(6.7))
bullets(s, Inches(7.6), Inches(1.7), Inches(5.2), Inches(4.8), [
    (0, [("Y = ", {}), ("max found %", {"bold": True}), (" (did you reach the tallest cliff).", {})]),
    (0, [("X = ", {}), ("coverage of q>5x", {"bold": True}), (" (how many cliffs you caught).", {})]),
    (0, [("ucb / ei / pi", {"font": FM, "color": ACC2, "bold": True}),
         (" occupy the high-coverage right side.", {})]),
    (0, [("sigma", {"font": FM, "color": ACCENT, "bold": True}),
         (" hits high peaks but sits at low coverage (explores broadly, not tail-focused).", {})]),
    (0, [("ts", {"font": FM, "bold": True}), (" is the randomised middle ground.", {})]),
])

# =========================================================
# 15 — CONCLUSION
# =========================================================
s = slide("Conclusion", "Axis-split BO wins — then pick the acquisition")
bullets(s, LX, Inches(1.5), Inches(12.23), Inches(2.5), [
    (0, [("Axis-split (one BO per axis) is the clear winner: ", {"bold": True, "color": GREEN}),
         ("point_split", {"font": FM, "bold": True}),
         (f" leads every representation mean — med-unseen q {by_rep.iloc[0].qmed:.2f}, "
          f"max-found {by_rep.iloc[0].maxf:.0f}%, coverage>5x {by_rep.iloc[0].cov5:.0f}% "
          f"(vs pair {by_rep.iloc[1].qmed:.2f} / point {by_rep.iloc[2].qmed:.2f}).", {})]),
    (0, [("Best boundary discovery: ", {"bold": True, "color": ACC2}),
         (f"{best_disc.method}", {"font": FM, "bold": True}),
         (f"  — {best_disc.maxf:.0f}% of the true max q-error, {best_disc.cov5:.0f}% coverage>5x "
          f"at {budget_pct:.0f}% budget.", {})]),
    (0, [("Best accuracy: ", {"bold": True, "color": ACCENT}),
         (f"{best_bo.method}", {"font": FM, "bold": True}),
         (f"  (med-unseen q {best_bo.qmed:.3f}) — sigma ties the linear baselines AND still finds boundaries.", {})]),
])
chip(s, LX, Inches(4.3), Inches(6.0), [("Hunt boundaries  →  ", {"bold": True}),
     ("point_split + ucb", {"font": FM})], ACC2, size=14, h=0.6)
chip(s, Inches(6.75), Inches(4.3), Inches(6.0), [("Accurate + finds cliffs  →  ", {"bold": True}),
     ("point_split + sigma", {"font": FM})], GREEN, size=14, h=0.6)
textbox(s, LX, Inches(5.25), Inches(12.23), Inches(1.6), [
    [("Why axis-split?", {"bold": True}), ("  the input stays exactly P, but running one BO per axis removes the "
      "horizontal/vertical X-collision — recovering the axis signal with NO extra feature.  ", {}),
     ("point_split is the recommended representation", {"bold": True, "color": GREEN}),
     (": sigma for an accurate surface, ucb for maximum boundary coverage — at a 10% budget a blind grid never targets.", {})],
], size=13.5, space_after=6)

# =========================================================
prs.save(str(OUT_PPTX))
print(f"saved: {OUT_PPTX}  ({len(prs.slides._sldIdLst)} slides)")
