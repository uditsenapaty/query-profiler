#!/usr/bin/env python3
"""Approximate visual preview of the pptx (no Office renderer available).
Exact shape geometry/colours from python-pptx; text approximated with wrapping.
Renders PNGs + a contact sheet so layout can be eyeballed via the Read tool."""
import sys, os, textwrap
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
EMU = 914400.0
def I(v): return (v or 0)/EMU

SRC = sys.argv[1] if len(sys.argv) > 1 else "../interpolation_methods_explained.pptx"
OUTDIR = sys.argv[2] if len(sys.argv)>2 else "../preview"
os.makedirs(OUTDIR, exist_ok=True)
prs = Presentation(SRC)
SW, SH = I(prs.slide_width), I(prs.slide_height)

def hexof(color):
    try:
        return "#%02x%02x%02x" % (color[0], color[1], color[2])
    except Exception:
        return None

def fill_hex(shape):
    try:
        if shape.fill.type is not None and shape.fill.type == 1:  # solid
            c = shape.fill.fore_color.rgb
            return "#%02x%02x%02x" % (c[0], c[1], c[2])
    except Exception:
        pass
    return None

def run_style(run):
    f = run.font
    sz = f.size.pt if f.size else 14
    col = "#1f2a37"
    try:
        if f.color and f.color.type is not None and f.color.rgb is not None:
            c = f.color.rgb; col = "#%02x%02x%02x" % (c[0], c[1], c[2])
    except Exception:
        pass
    mono = (f.name == "Consolas")
    return sz, col, bool(f.bold), mono

def draw_text_frame(ax, shape):
    l, t, w, h = I(shape.left), I(shape.top), I(shape.width), I(shape.height)
    y = t + 0.06
    for para in shape.text_frame.paragraphs:
        runs = [r for r in para.runs if r.text]
        if not runs:
            y += 0.16; continue
        sz, col, bold, mono = run_style(runs[0])
        txt = "".join(r.text for r in runs)
        cpl = max(6, int(w / (0.0095*sz if mono else 0.0072*sz)))
        wrapped = textwrap.wrap(txt, cpl) or [""]
        align = para.alignment
        for line in wrapped:
            if align == 2:      # right
                xx, ha = l + w - 0.06, "right"
            elif align == 3:    # center handled crudely as left
                xx, ha = l + w/2, "center"
            else:
                xx, ha = l + 0.08, "left"
            ax.text(xx, y + sz/72.0*0.5, line, fontsize=sz, color=col,
                    ha=ha, va="center",
                    family="monospace" if mono else "sans-serif",
                    fontweight="bold" if bold else "normal")
            y += sz/72.0*1.25
        y += 0.02

def draw_table(ax, shape):
    tbl = shape.table
    l, t, w, h = I(shape.left), I(shape.top), I(shape.width), I(shape.height)
    ncol = len(tbl.columns); nrow = len(tbl.rows)
    cw = [I(c.width) for c in tbl.columns]; tot = sum(cw) or w
    cw = [c/tot*w for c in cw]
    rh = [I(r.height) for r in tbl.rows]; rtot = sum(rh) or h
    rh = [r/rtot*h for r in rh]
    yy = t
    for i in range(nrow):
        xx = l
        for j in range(ncol):
            cell = tbl.cell(i, j)
            fh = None
            try:
                if cell.fill.type == 1:
                    c = cell.fill.fore_color.rgb; fh = "#%02x%02x%02x" % (c[0], c[1], c[2])
            except Exception:
                pass
            ax.add_patch(Rectangle((xx, yy), cw[j], rh[i], facecolor=fh or "white",
                                   edgecolor="#cfd6db", linewidth=0.4))
            runs = [r for p in cell.text_frame.paragraphs for r in p.runs if r.text]
            if runs:
                sz, col, bold, mono = run_style(runs[0])
                sz = min(sz, 11)
                txt = "".join(r.text for r in runs)
                al = cell.text_frame.paragraphs[0].alignment
                if al == 2: hx, ha = xx+cw[j]-0.04, "right"
                elif al is None or al == 1: hx, ha = xx+cw[j]/2, "center"
                else: hx, ha = xx+0.05, "left"
                ax.text(hx, yy+rh[i]/2, txt[:int(cw[j]/(0.011*sz))+2], fontsize=sz,
                        color=col, ha=ha, va="center",
                        family="monospace" if mono else "sans-serif",
                        fontweight="bold" if bold else "normal")
            xx += cw[j]
        yy += rh[i]

def render(slide, idx):
    fig = plt.figure(figsize=(SW, SH), dpi=110)
    ax = fig.add_axes([0, 0, 1, 1]); ax.set_xlim(0, SW); ax.set_ylim(0, SH)
    ax.invert_yaxis(); ax.axis("off")
    ax.add_patch(Rectangle((0, 0), SW, SH, facecolor="white", edgecolor="none"))
    for shape in slide.shapes:
        try:
            l, t, w, h = I(shape.left), I(shape.top), I(shape.width), I(shape.height)
        except Exception:
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            fh = fill_hex(shape)
            rounded = "ROUNDED" in str(shape.adjustments._sppr.getparent().tag).upper() if False else True
            ax.add_patch(FancyBboxPatch((l, t), w, h,
                         boxstyle="round,pad=0,rounding_size=0.04",
                         facecolor=fh or "none", edgecolor="#e8eced", linewidth=0.5))
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                import io
                img = plt.imread(io.BytesIO(shape.image.blob))
                ax.imshow(img, extent=(l, l+w, t+h, t), aspect="auto", zorder=3)
            except Exception as e:
                ax.add_patch(Rectangle((l,t),w,h,facecolor="#eef",edgecolor="#99c"))
        if shape.has_table:
            draw_table(ax, shape)
        elif shape.has_text_frame and shape.text_frame.text.strip():
            draw_text_frame(ax, shape)
    p = os.path.join(OUTDIR, f"slide_{idx:02d}.png")
    fig.savefig(p); plt.close(fig)
    return p

paths = [render(s, i) for i, s in enumerate(prs.slides, 1)]

# contact sheet
import math
n = len(paths); cols = 5; rows = math.ceil(n/cols)
fig, axs = plt.subplots(rows, cols, figsize=(cols*2.7, rows*1.6), dpi=110)
for k, ax in enumerate(axs.ravel()):
    ax.axis("off")
    if k < n:
        ax.imshow(plt.imread(paths[k])); ax.set_title(f"{k+1}", fontsize=7)
fig.tight_layout()
fig.savefig(os.path.join(OUTDIR, "contact_sheet.png")); plt.close(fig)
print("rendered", n, "slides to", os.path.abspath(OUTDIR))
