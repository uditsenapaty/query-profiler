#!/usr/bin/env python3
# =========================================================
# scripts/make_boundary_ppt.py
# ---------------------------------------------------------
# Focused deck explaining ONLY the grid-boundary adjustment
# in build_gt.py :: find_min_positive_point_nd() —
# choosing the minimum value per axis so the JOINT predicate
# cardinality is never 0 at the grid's lower boundary.
#
# 2D selectivity-space visuals (diagonal doubling + binary
# search + coordinate descent), old-vs-new range on a
# 10x10-resolution (11x11-point) grid, and a selectivity/
# value comparison.
#
# Run:  python3 scripts/make_boundary_ppt.py
# =========================================================

import os
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import Rectangle, FancyArrowPatch, Polygon

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------
# palette (matches the interpolation deck)
# ---------------------------------------------------------
INK    = RGBColor(0x1F, 0x2A, 0x37)
ACCENT = RGBColor(0x0E, 0x7C, 0x86)
ACC2   = RGBColor(0xB4, 0x53, 0x09)
GREEN  = RGBColor(0x15, 0x6F, 0x3B)
RED    = RGBColor(0xB1, 0x2A, 0x2A)
GREY   = RGBColor(0x6B, 0x72, 0x80)
LIGHT  = RGBColor(0xF2, 0xF4, 0xF6)
ROW2   = RGBColor(0xEC, 0xF3, 0xF4)
HIROW  = RGBColor(0xFD, 0xF1, 0xDD)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BAND   = RGBColor(0xE8, 0xEC, 0xEF)
ROW1   = WHITE
FT, FM = "Segoe UI", "Consolas"

# hex for matplotlib
H_DEAD = "#f6d6d3"; H_POS = "#dcebd6"
H_DEADL = "#c0392b"; H_POSL = "#2e7d32"
H_ACC = "#0e7c86"; H_ACC2 = "#b45309"; H_INK = "#1f2a37"; H_GREY = "#6b7280"

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
BLANK = prs.slide_layouts[6]
_page = [0]

# =========================================================
# model of the 2D selectivity space (illustrative)
#   joint card > 0  iff  s1 >= A  AND  s2 >= B
#   (real frontiers curve; per-axis floors are what the
#    algorithm returns, so a rectangular corner is faithful)
# =========================================================
A, B = 0.18, 0.30          # true per-axis selectivity floors
RES, NPTS = 10, 11         # 10x10 resolution -> 11x11 points

def card(s):               # joint surviving cardinality (>0 in positive region)
    return 1 if (s[0] >= A and s[1] >= B) else 0

# ---------------------------------------------------------
# pptx helpers (subset, same style as interpolation deck)
# ---------------------------------------------------------
def _runs(p, runs, size, color, font=FT, align=None):
    if align is not None: p.alignment = align
    if isinstance(runs, str): runs = [(runs, {})]
    for text, o in runs:
        r = p.add_run(); r.text = text; f = r.font
        f.name = o.get("font", font); f.size = Pt(o.get("size", size))
        f.bold = o.get("bold", False); f.italic = o.get("italic", False)
        f.color.rgb = o.get("color", color)
    return p

def textbox(s, l, t, w, h, lines, size=16, color=INK, font=FT,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=4, line_spacing=1.0):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(1)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after); p.space_before = Pt(0); p.line_spacing = line_spacing
        _runs(p, ln, size, color, font=font, align=align)
    return tb

def rect(s, l, t, w, h, fill):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.fill.background()
    sp.shadow.inherit = False; return sp

def rrect(s, l, t, w, h, fill, line=None, lw=1.0):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sp.adjustments[0] = 0.06; sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    sp.shadow.inherit = False; return sp

def slide(kicker, title, accent=ACCENT):
    s = prs.slides.add_slide(BLANK); _page[0] += 1
    rect(s, 0, 0, SW, Inches(0.16), accent)
    textbox(s, Inches(0.55), Inches(0.24), Inches(11.8), Inches(0.3),
            [[(kicker.upper(), {"size": 12.5, "bold": True, "color": accent})]])
    textbox(s, Inches(0.55), Inches(0.52), Inches(12.3), Inches(0.8),
            [[(title, {"size": 26, "bold": True, "color": INK})]])
    rect(s, Inches(0.55), Inches(1.26), Inches(12.23), Pt(1.6), BAND)
    textbox(s, Inches(0.55), Inches(7.12), Inches(9), Inches(0.3),
            [[("Grid boundary anchoring  ·  build_gt.py :: find_min_positive_point_nd", {"size": 9, "color": GREY})]])
    textbox(s, Inches(12.2), Inches(7.12), Inches(0.9), Inches(0.3),
            [[(str(_page[0]), {"size": 9, "color": GREY})]], align=PP_ALIGN.RIGHT)
    return s

def mono(s, l, t, w, h, lines, size=14, title=None, fill=LIGHT, accent=ACCENT, line_spacing=1.05):
    rrect(s, l, t, w, h, fill, line=BAND, lw=1.0)
    pad = Inches(0.14); ty = t + Inches(0.06)
    if title:
        textbox(s, l+pad, ty, w-2*pad, Inches(0.3),
                [[(title, {"size": size-1.5, "bold": True, "color": accent})]]); ty += Inches(0.34)
    textbox(s, l+pad, ty, w-2*pad, h-(ty-t)-Inches(0.06), lines, size=size, color=INK,
            font=FM, line_spacing=line_spacing, space_after=3)

def bullets(s, l, t, w, h, items, size=15, color=INK, accent=ACCENT, gap=5, line_spacing=1.03):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    marks = {0: "▸  ", 1: "–  ", 2: "·  "}
    for i, (lvl, runs) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl; p.space_after = Pt(gap); p.space_before = Pt(0); p.line_spacing = line_spacing
        mcol = accent if lvl == 0 else (GREY if lvl == 2 else ACCENT)
        r2 = [(marks.get(lvl, ""), {"color": mcol, "bold": True})]
        r2 += [(runs, {})] if isinstance(runs, str) else list(runs)
        _runs(p, r2, size if lvl == 0 else size-1.5, color)
    return tb

def _cell(cell, runs, size, color, bold, align, fill):
    cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.fill.solid(); cell.fill.fore_color.rgb = fill
    p = cell.text_frame.paragraphs[0]; p.alignment = align
    cell.text_frame.word_wrap = True
    if isinstance(runs, str): runs = [(runs, {})]
    for text, o in runs:
        r = p.add_run(); r.text = text; f = r.font
        f.name = o.get("font", FT); f.size = Pt(o.get("size", size))
        f.bold = o.get("bold", bold); f.italic = o.get("italic", False)
        f.color.rgb = o.get("color", color)

def table(s, l, t, w, header, rows, col_w=None, size=12, header_fill=ACCENT,
          hi=None, aligns=None, row_h=0.36, header_h=0.38):
    nC = len(header); nR = len(rows)+1
    h = Inches(header_h + row_h*len(rows))
    gt = s.shapes.add_table(nR, nC, l, t, w, h).table
    gt.first_row = False; gt.horz_banding = False
    gt._tbl.tblPr.set('firstRow', '0'); gt._tbl.tblPr.set('bandRow', '0')
    if col_w:
        tot = sum(col_w)
        for j, cwj in enumerate(col_w): gt.columns[j].width = Emu(int(w*cwj/tot))
    gt.rows[0].height = Inches(header_h)
    if aligns is None: aligns = [PP_ALIGN.LEFT]+[PP_ALIGN.CENTER]*(nC-1)
    for j, htx in enumerate(header): _cell(gt.cell(0, j), htx, size, WHITE, True, aligns[j], header_fill)
    hi = hi or set()
    for i, row in enumerate(rows):
        gt.rows[i+1].height = Inches(row_h)
        base = HIROW if i in hi else (ROW2 if i % 2 else ROW1)
        for j, c in enumerate(row): _cell(gt.cell(i+1, j), c, size, INK, False, aligns[j], base)
    return gt

def chip(s, l, t, w, text, fill, tcolor=WHITE, size=12.5, h=0.42):
    sp = rrect(s, l, t, w, Inches(h), fill); tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.08); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    _runs(p, text if isinstance(text, list) else [(text, {})], size, tcolor)
    return sp

def add_image(s, path, l, t, w=None, h=None):
    return s.shapes.add_picture(path, l, t,
                                width=Inches(w) if w else None,
                                height=Inches(h) if h else None)

# =========================================================
# matplotlib figures
# =========================================================
FIGDIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "boundary_fig")
os.makedirs(FIGDIR, exist_ok=True)

def _base(ax, grid=None, title=None):
    ax.add_patch(Rectangle((0, 0), 1, 1, fc=H_DEAD, ec="none"))             # dead (default)
    ax.add_patch(Rectangle((A, B), 1-A, 1-B, fc=H_POS, ec="none"))          # positive corner
    ax.plot([A, A], [B, 1], color=H_POSL, lw=1.5)
    ax.plot([A, 1], [B, B], color=H_POSL, lw=1.5)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.set_aspect("equal")
    ax.set_xlabel("selectivity  s₁  (axis x1)", fontsize=12)
    ax.set_ylabel("selectivity  s₂  (axis x2)", fontsize=12)
    ax.tick_params(labelsize=9)
    if title: ax.set_title(title, fontsize=13, color=H_INK, weight="bold")

def _legend_text(ax):
    ax.text(0.045, 0.045, "joint card = 0\n(empty result)", fontsize=10.5,
            color=H_DEADL, weight="bold", va="bottom")
    ax.text(0.97, 0.965, "joint card > 0", fontsize=10.5, color=H_POSL,
            weight="bold", ha="right", va="top")

def fig_problem():
    fig, ax = plt.subplots(figsize=(6.0, 5.4), dpi=130)
    _base(ax, title="OLD range: grid spans the full axis (global min → max)")
    g = np.linspace(0, 1, NPTS)
    for x in g:
        for y in g:
            if card([x, y]) > 0:
                ax.plot(x, y, "o", ms=5, color=H_POSL)
            else:
                ax.plot(x, y, "x", ms=6, color=H_DEADL, mew=1.6)
    n_dead = sum(1 for x in g for y in g if card([x, y]) == 0)
    _legend_text(ax)
    ax.text(0.5, -0.16, f"{n_dead} of {NPTS*NPTS} grid points are DEAD (0 rows)",
            transform=ax.transAxes, ha="center", fontsize=12, color=H_DEADL, weight="bold")
    fig.tight_layout(); p = os.path.join(FIGDIR, "problem.png"); fig.savefig(p, bbox_inches="tight"); plt.close(fig)
    return p, n_dead

def fig_diagonal():
    fig, ax = plt.subplots(figsize=(6.0, 5.4), dpi=130)
    _base(ax, title="Phase 1–2: diagonal doubling, then binary search")
    ax.plot([0, 1], [0, 1], "--", color=H_GREY, lw=1.0)
    # phase 1: doubling along diagonal (schematic; real eps=1e-9)
    doubles = [0.04, 0.08, 0.16, 0.32]
    for t in doubles:
        pos = card([t, t]) > 0
        ax.plot(t, t, "o", ms=8, mfc=(H_POS if pos else "white"),
                mec=(H_POSL if pos else H_DEADL), mew=1.8)
    for i in range(len(doubles)-1):
        ax.annotate("", xy=(doubles[i+1], doubles[i+1]), xytext=(doubles[i], doubles[i]),
                    arrowprops=dict(arrowstyle="->", color=H_ACC2, lw=1.6))
    ax.text(0.05, 0.012, "×2 doublings (eps=1e-9 …)", color=H_ACC2, fontsize=10, weight="bold")
    # phase 2: binary search bracket on diagonal between 0.16 (dead) and 0.32 (pos)
    for mid in [0.24, 0.28, 0.30]:
        ax.plot(mid, mid, "s", ms=6, color=H_ACC, alpha=0.55)
    ax.plot(B, B, "*", ms=20, color=H_ACC2, mec="black", mew=0.5, zorder=6)
    ax.annotate("s* = (0.30, 0.30)\nsmallest diagonal\npoint with card>0",
                xy=(B, B), xytext=(0.5, 0.16), fontsize=10.5, color=H_INK,
                arrowprops=dict(arrowstyle="->", color="black", lw=1.2))
    fig.tight_layout(); p = os.path.join(FIGDIR, "diagonal.png"); fig.savefig(p, bbox_inches="tight"); plt.close(fig)
    return p

def fig_coord():
    fig, ax = plt.subplots(figsize=(6.0, 5.4), dpi=130)
    _base(ax, title="Phase 3: coordinate descent → true per-axis floors")
    # conservative square from diagonal
    ax.add_patch(Rectangle((B, B), 1-B, 1-B, fc="none", ec=H_GREY, ls="--", lw=1.3))
    ax.plot(B, B, "*", ms=18, color=H_ACC2, mec="black", mew=0.5, zorder=6)
    ax.text(B+0.015, B-0.045, "diagonal box (0.30,0.30)", fontsize=9.5, color=H_GREY)
    # descend s1: 0.30 -> 0.18 (s2 held at 0.30)
    ax.annotate("", xy=(A, B), xytext=(B, B),
                arrowprops=dict(arrowstyle="->", color=H_ACC, lw=2.4))
    ax.text((A+B)/2, B+0.02, "lower s₁\n0.30→0.18", ha="center", fontsize=10, color=H_ACC, weight="bold")
    # attempt s2: held (already minimal at 0.30)
    ax.annotate("", xy=(A, B-0.0), xytext=(A, B),
                arrowprops=dict(arrowstyle="-", color=H_GREY, lw=1.0))
    ax.text(A-0.015, (B+0.05), "s₂ already\nminimal (0.30)", ha="right", fontsize=9.5, color=H_GREY)
    # final corner
    ax.add_patch(Rectangle((A, B), 1-A, 1-B, fc="none", ec=H_POSL, lw=2.2))
    ax.plot(A, B, "o", ms=11, color=H_POSL, mec="black", mew=0.6, zorder=7)
    ax.annotate("anchor = (0.18, 0.30)\nper-axis minimum sels\nwith joint card>0",
                xy=(A, B), xytext=(0.42, 0.62), fontsize=10.5, color=H_POSL, weight="bold",
                arrowprops=dict(arrowstyle="->", color=H_POSL, lw=1.4))
    fig.tight_layout(); p = os.path.join(FIGDIR, "coord.png"); fig.savefig(p, bbox_inches="tight"); plt.close(fig)
    return p

def fig_oldnew():
    fig, axs = plt.subplots(1, 2, figsize=(11.6, 5.0), dpi=130)
    g_old = np.linspace(0, 1, NPTS)
    g1 = np.linspace(A, 1, NPTS); g2 = np.linspace(B, 1, NPTS)
    # OLD
    ax = axs[0]; _base(ax, title="OLD range  [global min, max]")
    nd = 0
    for x in g_old:
        for y in g_old:
            if card([x, y]) > 0: ax.plot(x, y, "o", ms=4.5, color=H_POSL)
            else: ax.plot(x, y, "x", ms=5.5, color=H_DEADL, mew=1.5); nd += 1
    ax.text(0.5, -0.155, f"{nd}/{NPTS*NPTS} points DEAD", transform=ax.transAxes,
            ha="center", fontsize=12, color=H_DEADL, weight="bold")
    # NEW
    ax = axs[1]; _base(ax, title="NEW range  [anchor, max]")
    for x in g1:
        for y in g2:
            ax.plot(x, y, "o", ms=4.5, color=H_POSL)
    ax.add_patch(Rectangle((A, B), 1-A, 1-B, fc="none", ec=H_POSL, lw=2.0))
    ax.plot(A, B, "o", ms=10, color=H_ACC2, mec="black", mew=0.6, zorder=7)
    ax.text(0.5, -0.155, f"0/{NPTS*NPTS} dead — every point has rows",
            transform=ax.transAxes, ha="center", fontsize=12, color=H_POSL, weight="bold")
    fig.tight_layout(); p = os.path.join(FIGDIR, "oldnew.png"); fig.savefig(p, bbox_inches="tight"); plt.close(fig)
    return p, nd

def _axstyle(ax):
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.set_aspect("equal")
    ax.set_xlabel("selectivity  s₁", fontsize=11); ax.set_ylabel("selectivity  s₂", fontsize=11)
    ax.tick_params(labelsize=8)

def fig_staircase():
    fig, axs = plt.subplots(1, 2, figsize=(11.4, 5.2), dpi=130)
    # ---- panel A: clean rectangular corner ----
    ax = axs[0]; a, b = 0.25, 0.25
    ax.add_patch(Rectangle((0, 0), 1, 1, fc=H_DEAD))
    ax.add_patch(Rectangle((a, b), 1-a, 1-b, fc=H_POS))
    ax.plot([a, a], [b, 1], color=H_POSL, lw=1.5); ax.plot([a, 1], [b, b], color=H_POSL, lw=1.5)
    ax.add_patch(Rectangle((a, b), 1-a, 1-b, fc="none", ec=H_POSL, lw=2.2))
    for x in np.linspace(a, 1, 11):
        for y in np.linspace(b, 1, 11): ax.plot(x, y, "o", ms=2.4, color=H_POSL)
    ax.plot(a, b, "o", ms=11, color=H_ACC2, mec="black", mew=0.6, zorder=6)
    ax.set_title("fig 1 · rectangular corner  →  nothing lost", fontsize=12.5, weight="bold", color=H_INK)
    ax.text(0.5, -0.135, "inscribed rectangle = the whole positive region",
            transform=ax.transAxes, ha="center", fontsize=11, color=H_POSL, weight="bold")
    _axstyle(ax)
    # ---- panel B: staircase corner ----
    ax = axs[1]
    stair = [(0.2, 1.0), (0.2, 0.6), (0.4, 0.6), (0.4, 0.4), (0.6, 0.4), (0.6, 0.3), (1.0, 0.3)]
    ax.add_patch(Rectangle((0, 0), 1, 1, fc=H_POS))
    dead = [(0, 0), (0, 1), (0.2, 1), (0.2, 0.6), (0.4, 0.6), (0.4, 0.4),
            (0.6, 0.4), (0.6, 0.3), (1.0, 0.3), (1.0, 0)]
    ax.add_patch(Polygon(dead, closed=True, fc=H_DEAD, ec="none"))
    ax.plot(*zip(*stair), color=H_POSL, lw=1.8)
    cx, cy = 0.4, 0.4
    ax.add_patch(Rectangle((cx, cy), 1-cx, 1-cy, fc="none", ec=H_POSL, lw=2.4))
    # lost positive arms (off the grid rectangle)
    for arm in ([(0.2, 0.6), (0.2, 1.0), (0.4, 1.0), (0.4, 0.6)],
                [(0.6, 0.3), (0.6, 0.4), (1.0, 0.4), (1.0, 0.3)]):
        ax.add_patch(Polygon(arm, closed=True, fc="#f3c98f", ec=H_ACC2, lw=1.4, hatch="//"))
    for x, y in [(0.30, 0.86), (0.30, 0.70), (0.80, 0.35), (0.92, 0.35)]:
        ax.plot(x, y, "*", ms=13, color=H_ACC2, mec="black", mew=0.4, zorder=7)
    ax.plot(cx, cy, "o", ms=11, color=H_ACC2, mec="black", mew=0.6, zorder=6)
    ax.text(cx+0.02, cy+0.04, "grid = [0.4,1]²", fontsize=10, color=H_POSL, weight="bold")
    ax.annotate("lost — still positive,\nbut off the grid", xy=(0.30, 0.78), xytext=(0.005, 0.46),
                fontsize=9.5, color=H_ACC2, weight="bold",
                arrowprops=dict(arrowstyle="->", color=H_ACC2, lw=1.2))
    ax.annotate("", xy=(0.86, 0.35), xytext=(0.62, 0.12),
                arrowprops=dict(arrowstyle="->", color=H_ACC2, lw=1.2))
    ax.text(0.45, 0.075, "lost", fontsize=9.5, color=H_ACC2, weight="bold")
    ax.set_title("fig 2 · staircase corner  →  loses near-frontier points",
                 fontsize=12.5, weight="bold", color=H_INK)
    ax.text(0.5, -0.135, "one rectangle can't hug every step",
            transform=ax.transAxes, ha="center", fontsize=11, color=H_ACC2, weight="bold")
    _axstyle(ax)
    fig.tight_layout(); p = os.path.join(FIGDIR, "limitation.png"); fig.savefig(p, bbox_inches="tight"); plt.close(fig)
    return p

P_stair = fig_staircase()
P_problem, N_DEAD = fig_problem()
P_diag = fig_diagonal()
P_coord = fig_coord()
P_oldnew, N_DEAD2 = fig_oldnew()

# =========================================================
# SLIDES
# =========================================================
LX = Inches(0.55)

# 1 — title
s = prs.slides.add_slide(BLANK); _page[0] += 1
rect(s, 0, 0, SW, SH, INK)
rect(s, 0, Inches(4.5), SW, Inches(0.05), ACCENT)
textbox(s, Inches(0.8), Inches(1.4), Inches(11.7), Inches(2.6), [
    [("Anchoring the Grid", {"size": 40, "bold": True, "color": WHITE})],
    [("Choosing each axis's minimum so the joint cardinality is never 0 at the boundary",
      {"size": 18, "color": RGBColor(0xC8, 0xD2, 0xDA)})],
], space_after=10)
textbox(s, Inches(0.8), Inches(4.75), Inches(11.7), Inches(1.6), [
    [("build_gt.py", {"size": 16, "font": FM, "color": RGBColor(0x6F, 0xD3, 0xDA)}),
     (" :: ", {"size": 16, "color": RGBColor(0x9F, 0xB0, 0xBC)}),
     ("find_min_positive_point_nd()", {"size": 16, "font": FM, "color": RGBColor(0x6F, 0xD3, 0xDA)})],
    [("diagonal doubling · binary search · coordinate descent · 10×10 resolution (11×11 points)",
      {"size": 14, "color": RGBColor(0x9F, 0xB0, 0xBC)})],
], space_after=12)

# 2 — the problem
s = slide("The problem", "Starting at the global min → an empty lower corner")
bullets(s, LX, Inches(1.5), Inches(6.0), Inches(3.6), [
    (0, [("The grid is built over each axis's value range.  The naive choice is ", {}),
         ("[global min, max]", {"font": FM, "bold": True}), (".", {})]),
    (0, [("At low selectivity on ", {}), ("all", {"italic": True}),
         (" axes, the joint predicate survives ", {}), ("0 rows", {"bold": True, "color": RED}), (".", {})]),
    (0, "That whole lower-left strip is wasted:"),
    (1, [("empty COUNT(*) → trivial/degenerate plan & runtime,", {})]),
    (1, [("meaningless q-errors, budget spent on dead points.", {})]),
    (0, [("Goal: ", {"bold": True}), ("guarantee joint card > 0 at the lower boundary.", {"bold": True, "color": GREEN})]),
])
add_image(s, P_problem, Inches(7.0), Inches(1.4), h=5.5)

# 3 — the fix / card_fn + plan
s = slide("The fix", "Search selectivity space for the per-axis minimum with rows")
mono(s, LX, Inches(1.5), Inches(6.05), Inches(2.0), [
    "card_fn(s1..sd):",
    "   vᵢ = sorted_valsᵢ[ ⌊ sᵢ·(nᵢ−1) ⌋ ]   # sel → value",
    "   return COUNT(*) of joint predicate(v1..vd)",
], size=12.5, title="map selectivities → values → joint COUNT(*)")
bullets(s, LX, Inches(3.7), Inches(6.05), Inches(3.1), [
    (0, [("Find the smallest selectivities ", {}), ("(s1..sd)", {"font": FM}),
         (" with ", {}), ("card_fn > 0", {"font": FM, "bold": True}), (", in 3 phases:", {})]),
    (1, [("diagonal expansion", {"bold": True, "color": ACCENT}), (" — leap onto the positive region,", {})]),
    (1, [("binary search", {"bold": True, "color": ACCENT}), (" — pin the diagonal crossing,", {})]),
    (1, [("coordinate descent", {"bold": True, "color": ACCENT}), (" — minimise each axis on its own.", {})]),
])
chip(s, Inches(7.05), Inches(1.5), Inches(5.75),
     [("monotone assumption: more selectivity ⇒ ≥ rows", {"bold": True})], ACCENT, size=12.5, h=0.5)
table(s, Inches(7.05), Inches(2.2), Inches(5.75),
      ["build_gt.py output", "meaning"],
      [
       [[("boundary_sels", {"font": FM, "size": 11})], "per-axis min selectivities"],
       [[("boundary_count", {"font": FM, "size": 11})], "joint rows at that anchor (>0)"],
       [[("axis_lower_bounds", {"font": FM, "size": 11})], "→ the value each grid starts at"],
      ], col_w=[2.5, 3.3], size=11.5, row_h=0.5)
textbox(s, Inches(7.05), Inches(4.5), Inches(5.75), Inches(2.0), [
    [("axis_lower_bounds", {"font": FM, "bold": True, "size": 12}),
     (" is passed to ", {"size": 12}), ("sampler.sample(lower_bound=…)", {"font": FM, "size": 11.5}),
     (", so each axis grid spans ", {"size": 12}),
     ("[anchor, max]", {"font": FM, "bold": True, "size": 12}), (" instead of ", {"size": 12}),
     ("[global min, max]", {"font": FM, "size": 12}), (".", {"size": 12})],
])

# 4 — diagonal + binary search
s = slide("Phases 1–2 · the diagonal", "Leap onto the region, then pin the crossing")
add_image(s, P_diag, Inches(0.5), Inches(1.4), h=5.5)
mono(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(2.05), [
    "lo,hi = 0, eps            # eps = 1e-9",
    "while card([hi]*d)==0:    # PHASE 1",
    "    lo, hi = hi, hi*2     #   ×2 leap",
    "while hi-lo > tol:        # PHASE 2",
    "    m=(lo+hi)/2",
    "    hi=m if card([m]*d)>0 else …; lo=…",
], size=11.5, title="diagonal: same selectivity on every axis  (s,s,…,s)")
bullets(s, Inches(7.0), Inches(3.75), Inches(5.85), Inches(3.0), [
    (0, [("Phase 1", {"bold": True, "color": ACC2}), (" doubles ", {}), ("s", {"font": FM}),
         (" from ~0 until the diagonal point enters the positive region (cheap, geometric).", {})]),
    (0, [("Phase 2", {"bold": True, "color": ACC2}), (" binary-searches that last interval → ", {}),
         ("s* ", {"font": FM, "bold": True}), ("= smallest equal-selectivity point with rows.", {})]),
    (0, [("Gives a safe but ", {}), ("conservative square", {"bold": True}), (" box (s*, s*).", {})]),
])

# 5 — coordinate descent
s = slide("Phase 3 · coordinate descent", "Free each axis down to its own true minimum")
add_image(s, P_coord, Inches(0.5), Inches(1.4), h=5.5)
mono(s, Inches(7.0), Inches(1.5), Inches(5.8), Inches(2.05), [
    "improved = True",
    "while improved:",
    "  for axis in range(d):",
    "    lo,hi = 0, s[axis]      # binary search",
    "    ...  hi=m if card(trial)>0 else lo=m",
    "    if hi < s[axis]: s[axis]=hi; improved=True",
], size=11.5, title="per-axis binary search, others held fixed")
bullets(s, Inches(7.0), Inches(3.75), Inches(5.85), Inches(3.0), [
    (0, [("The diagonal forced both axes equal (0.30).  Coordinate descent ", {}),
         ("frees each axis", {"bold": True}), (" to its real floor.", {})]),
    (0, [("Here s₁ drops ", {}), ("0.30 → 0.18", {"font": FM, "bold": True, "color": ACCENT}),
         ("; s₂ is already minimal (0.30).", {})]),
    (0, [("Result: the ", {}), ("anchor corner (0.18, 0.30)", {"bold": True, "color": GREEN}),
         (" — the tight lower-left of the positive region.", {})]),
])

# 6 — old vs new range (11x11)
s = slide("Old vs new range", "10×10 resolution (11×11 points): every point now has rows", accent=ACC2)
add_image(s, P_oldnew, Inches(0.7), Inches(1.5), w=12.0)
textbox(s, LX, Inches(6.55), Inches(12.23), Inches(0.6), [
    [("Same 121 points, same evenly-spaced selectivities — only the per-axis ", {"size": 13}),
     ("start", {"size": 13, "bold": True}), (" moves from the global min to the anchor.  ", {"size": 13}),
     (f"Dead points: {N_DEAD2} → 0.", {"size": 13, "bold": True, "color": GREEN})],
])

# 7 — selectivity & value comparison
s = slide("Selectivity & value comparison", "What the anchor changes, per axis", accent=ACC2)
table(s, LX, Inches(1.5), Inches(12.23),
      ["axis", "OLD lower selectivity", "OLD lower value", "NEW lower selectivity", "NEW lower value"],
      [
       [[("x1", {"font": FM, "bold": True})], "≈ 0.00  (global min)",
        [("vals[0]", {"font": FM})], [("0.18", {"bold": True, "color": ACC2})],
        [("vals[⌊0.18·(n−1)⌋]", {"font": FM, "size": 10.5})]],
       [[("x2", {"font": FM, "bold": True})], "≈ 0.00  (global min)",
        [("vals[0]", {"font": FM})], [("0.30", {"bold": True, "color": ACC2})],
        [("vals[⌊0.30·(n−1)⌋]", {"font": FM, "size": 10.5})]],
      ], col_w=[1.0, 3.0, 2.4, 3.0, 2.9], size=12, header_fill=ACC2, row_h=0.55)
table(s, LX, Inches(3.3), Inches(12.23),
      ["at the lower-left corner", "OLD range", "NEW range"],
      [
       [[("joint selectivities", {})], [("(≈0, ≈0)", {"font": FM})], [("(0.18, 0.30)", {"font": FM, "bold": True})]],
       [[("joint cardinality (rows)", {})], [("0   (empty)", {"font": FM, "color": RED})],
        [("boundary_count  > 0", {"font": FM, "color": GREEN, "bold": True})]],
       [[("dead grid points (of 121)", {})], [(f"{N_DEAD}", {"font": FM, "color": RED, "bold": True})],
        [("0", {"font": FM, "color": GREEN, "bold": True})]],
      ], col_w=[3.4, 2.9, 3.0], size=12, header_fill=INK, row_h=0.46, hi={1})
textbox(s, LX, Inches(5.7), Inches(12.23), Inches(1.1), [
    [("value at a selectivity:  ", {"size": 12.5}),
     ("vals[ ⌊ sel · (n−1) ⌋ ]", {"font": FM, "bold": True, "size": 13}),
     ("   over that axis's sorted distinct values (n of them).", {"size": 12.5})],
    [("The selectivity grid is unchanged in spacing — it is just shifted so its origin sits on the anchor, not on the empty global min.",
      {"size": 12, "italic": True, "color": GREY})],
])

# 7b — m1 worked example: old vs new selectivity grid (real qt5 columns)
def _load_m1():
    anch = {"x1": 0.18, "x2": 0.30}
    try:
        import pandas as _pd
        base = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..")
        df = _pd.read_csv(os.path.join(base, "gt_results_sf1_10x10_s1q2/qt5/m1/ground_truth.csv"))
        out = {}
        for ax in ["x1", "x2"]:
            sub = df[[ax, f"selectivity_{ax}"]].dropna().drop_duplicates().sort_values(f"selectivity_{ax}")
            sr = sub[f"selectivity_{ax}"].values.astype(float); vr = sub[ax].values.astype(float)
            sold = np.linspace(0, 1, NPTS); vold = np.interp(sold, sr, vr)
            a = anch[ax]; snew = np.linspace(a, 1, NPTS); vnew = np.interp(snew, sr, vr)
            out[ax] = (sold, vold, snew, vnew, a)
        return out
    except Exception:                                   # fallback to precomputed
        return {
            "x1": (np.linspace(0, 1, 11),
                   np.array([-985.6, 107.89, 1227.72, 2303.04, 3392.18, 4485.36, 5599.46, 6705.35, 7801.94, 8895.4, 9999.99]),
                   np.linspace(0.18, 1, 11),
                   np.array([1003.76, 1894.42, 2782.26, 3676.41, 4574.49, 5488.05, 6395.70, 7297.51, 8195.59, 9094.22, 9999.99]), 0.18),
            "x2": (np.linspace(0, 1, 11),
                   np.array([-964.32, 142.10, 1245.79, 2352.88, 3446.26, 4558.44, 5616.01, 6727.61, 7807.14, 8897.56, 9999.72]),
                   np.linspace(0.30, 1, 11),
                   np.array([2352.88, 3118.25, 3891.13, 4664.20, 5404.50, 6171.81, 6943.52, 7699.19, 8461.39, 9228.21, 9999.72]), 0.30),
        }
M1 = _load_m1()

s = slide("Selectivity-space sampling (m1)", "m1 worked example — old vs new grid (real qt5 columns)", accent=ACC2)

def _m1_table(s, l, ax):
    sold, vold, snew, vnew, a = M1[ax]
    rows = []
    for i in range(NPTS):
        below = sold[i] < a - 1e-9
        rows.append([
            [(str(i), {"font": FM, "size": 10})],
            [(f"{sold[i]:.2f}", {"font": FM, "size": 10})],
            [(f"{vold[i]:,.0f}", {"font": FM, "size": 10, "color": RED if below else INK,
                                  "bold": below})],
            [(f"{snew[i]:.3f}", {"font": FM, "size": 10})],
            [(f"{vnew[i]:,.0f}", {"font": FM, "size": 10, "color": GREEN if i == 0 else INK,
                                  "bold": i == 0})],
        ])
    return table(s, l, Inches(2.0), Inches(5.95),
                 ["i", "old sel", "old val", "new sel", "new val"],
                 rows, col_w=[0.4, 0.95, 1.45, 0.95, 1.45], size=10,
                 header_fill=ACC2, row_h=0.30, hi={0})

chip(s, Inches(0.55), Inches(1.5), Inches(5.95),
     [("axis x1   ", {"bold": True}), ("anchor sel = 0.18", {"font": FM})], ACCENT, size=12.5, h=0.42)
_m1_table(s, Inches(0.55), "x1")
chip(s, Inches(6.85), Inches(1.5), Inches(5.95),
     [("axis x2   ", {"bold": True}), ("anchor sel = 0.30", {"font": FM})], ACCENT, size=12.5, h=0.42)
_m1_table(s, Inches(6.85), "x2")
textbox(s, LX, Inches(5.78), Inches(12.23), Inches(1.25), [
    [("m1 places points at ", {"size": 12.5}), ("evenly-spaced selectivities", {"size": 12.5, "bold": True}),
     (".  OLD spans [0,1] → its first points sit at the global min (", {"size": 12.5}),
     ("x1 = −986, x2 = −964", {"font": FM, "size": 12, "color": RED}),
     (").  NEW spans [anchor,1] → compressed above the anchor, so the lowest value sampled is the anchor's (", {"size": 12.5}),
     ("x1 ≈ 1004 @ .18,  x2 ≈ 2353 @ .30", {"font": FM, "size": 12, "color": GREEN}),
     (").", {"size": 12.5})],
    [("Same column distribution — only the sampled percentiles shift.  Real qt5 m1 columns; new values interpolated from the measured 11-pt curve; anchor (0.18, 0.30) is the running example (qt5's own corner already survives, so its true anchor ≈ 0).",
      {"size": 10.5, "italic": True, "color": GREY})],
], space_after=6)

# 7b-2 — m0 (data-space uniform) and m2 (Picasso geometric) old vs new
_RAW = {
    "x1": dict(
        m0_sel=np.array([0.0014, 0.1017, 0.1996, 0.3016, 0.4023, 0.5025, 0.6011, 0.7003, 0.8003, 0.9006, 1.0]),
        m0_val=np.array([-985.6, 113.0, 1211.5, 2310.1, 3408.6, 4507.2, 5605.8, 6704.3, 7802.9, 8901.4, 10000.0]),
        m2_f=np.array([0.0, 0.0004, 0.0013, 0.0033, 0.0077, 0.0176, 0.0398, 0.0893, 0.2, 0.4473, 1.0]),
        m2_val=np.array([-985.6, -981.1, -970.9, -948.8, -901.9, -795.2, -549.2, -7.4, 1227.7, 3909.9, 10000.0])),
    "x2": dict(
        m0_sel=np.array([0.003, 0.1013, 0.2009, 0.2999, 0.3992, 0.4978, 0.601, 0.6994, 0.8004, 0.9008, 1.0]),
        m0_val=np.array([-964.3, 132.1, 1228.5, 2324.9, 3421.3, 4517.7, 5614.1, 6710.5, 7806.9, 8903.3, 9999.7]),
        m2_f=np.array([0.0, 0.0004, 0.0013, 0.0033, 0.0077, 0.0176, 0.0398, 0.0893, 0.2, 0.4473, 1.0]),
        m2_val=np.array([-964.3, -957.7, -944.8, -925.0, -878.7, -785.2, -541.6, 20.7, 1245.8, 3955.3, 9999.7])),
}

def _load_grids():
    anch = {"x1": 0.18, "x2": 0.30}
    raw = _RAW
    try:
        import pandas as _pd
        base = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..")
        m0 = _pd.read_csv(os.path.join(base, "gt_results_sf1_10x10_s1q2/qt5/m0/ground_truth.csv"))
        m2 = _pd.read_csv(os.path.join(base, "gt_results_sf1_10x10_s1q2/qt5/m2/ground_truth.csv"))
        raw = {}
        for ax in ["x1", "x2"]:
            a0 = m0[[ax, f"selectivity_{ax}"]].dropna().drop_duplicates().sort_values(ax)
            a2 = m2[[ax, f"selectivity_{ax}"]].dropna().drop_duplicates().sort_values(f"selectivity_{ax}")
            raw[ax] = dict(m0_sel=a0[f"selectivity_{ax}"].values.astype(float),
                           m0_val=a0[ax].values.astype(float),
                           m2_f=a2[f"selectivity_{ax}"].values.astype(float),
                           m2_val=a2[ax].values.astype(float))
    except Exception:
        pass
    out = {"m0": {}, "m2": {}}
    for ax in ["x1", "x2"]:
        a = anch[ax]; R = raw[ax]; sel, val = R["m0_sel"], R["m0_val"]
        av = float(np.interp(a, sel, val)); mx = float(val.max())
        nv = np.linspace(av, mx, 11); nsl = np.interp(nv, val, sel)
        out["m0"][ax] = (sel, val, nsl, nv, a)
        f = R["m2_f"]; ov = R["m2_val"]; nsl2 = a + (1 - a) * f; nv2 = np.interp(nsl2, sel, val)
        out["m2"][ax] = (f, ov, nsl2, nv2, a)
    return out
GR = _load_grids()

def _compare_slide(kicker, title, key, sel_dec, note):
    s = slide(kicker, title, accent=ACC2)
    def tbl(l, ax):
        osl, ov, nsl, nv, a = GR[key][ax]
        rows = []
        for i in range(NPTS):
            below = osl[i] < a - 1e-9
            rows.append([
                [(str(i), {"font": FM, "size": 10})],
                [(f"{osl[i]:.{sel_dec}f}", {"font": FM, "size": 10})],
                [(f"{ov[i]:,.0f}", {"font": FM, "size": 10, "color": RED if below else INK, "bold": below})],
                [(f"{nsl[i]:.{sel_dec}f}", {"font": FM, "size": 10})],
                [(f"{nv[i]:,.0f}", {"font": FM, "size": 10, "color": GREEN if i == 0 else INK, "bold": i == 0})],
            ])
        return table(s, l, Inches(2.0), Inches(5.95), ["i", "old sel", "old val", "new sel", "new val"],
                     rows, col_w=[0.4, 0.95, 1.45, 0.95, 1.45], size=10, header_fill=ACC2, row_h=0.30, hi={0})
    chip(s, Inches(0.55), Inches(1.5), Inches(5.95),
         [("axis x1   ", {"bold": True}), (f"anchor sel = {GR[key]['x1'][4]}", {"font": FM})], ACCENT, size=12.5, h=0.42)
    tbl(Inches(0.55), "x1")
    chip(s, Inches(6.85), Inches(1.5), Inches(5.95),
         [("axis x2   ", {"bold": True}), (f"anchor sel = {GR[key]['x2'][4]}", {"font": FM})], ACCENT, size=12.5, h=0.42)
    tbl(Inches(6.85), "x2")
    textbox(s, LX, Inches(5.78), Inches(12.23), Inches(1.25), note, space_after=6)
    return s

_compare_slide("Data-space sampling (m0)", "m0 worked example — old vs new grid (real qt5 columns)", "m0", 3, [
    [("m0 places points ", {"size": 12.5}), ("evenly in VALUE", {"size": 12.5, "bold": True}),
     (".  OLD spans [global min, max] → first points at the global min (", {"size": 12.5}),
     ("x1 = −986, x2 = −964", {"font": FM, "size": 12, "color": RED}),
     (").  NEW spans [anchor value, max] → the same even spacing, shifted up to start at the anchor (", {"size": 12.5}),
     ("x1 ≈ 992 @ .18,  x2 ≈ 2326 @ .30", {"font": FM, "size": 12, "color": GREEN}), (").", {"size": 12.5})],
    [("Real qt5 m0 columns; new selectivities read off the value↔selectivity curve; anchor (0.18, 0.30) is the running example.",
      {"size": 10.5, "italic": True, "color": GREY})],
])

_compare_slide("Selectivity-space sampling (m2 · Picasso)", "m2 worked example — old vs new grid (real qt5 columns)", "m2", 4, [
    [("m2 clusters points ", {"size": 12.5}), ("geometrically at low selectivity", {"size": 12.5, "bold": True}),
     (".  OLD wastes ", {"size": 12.5}), ("~8 of 11 points", {"size": 12.5, "bold": True, "color": RED}),
     (" in the near-min region (", {"size": 12.5}),
     ("x1 sels ≤ 0.09 → vals −985…−7", {"font": FM, "size": 12, "color": RED}),
     (").  NEW restricts the percentiles to [anchor, max], so the same geometric cluster sits ", {"size": 12.5}),
     ("just above the anchor", {"size": 12.5, "bold": True, "color": GREEN}),
     (" — the method that gains most from anchoring.", {"size": 12.5})],
    [("Actual global selectivities shown; m2's Picasso target fractions are unchanged, only now measured within [anchor, max]; new values interpolated.",
      {"size": 10.5, "italic": True, "color": GREY})],
])

# 7c — the one limitation: staircase frontier
s = slide("Limitation", "The one limitation — a staircase frontier", accent=ACC2)
add_image(s, P_stair, Inches(0.4), Inches(1.5), w=7.55)
bullets(s, Inches(8.15), Inches(1.55), Inches(4.65), Inches(5.3), [
    (0, [("The grid is a Cartesian product → an ", {}),
         ("axis-aligned rectangle", {"bold": True}), (" in selectivity space.", {})]),
    (0, [("Phase 3 finds ", {}), ("one", {"bold": True, "color": ACC2}),
         (" corner and inscribes ", {}), ("one", {"bold": True}), (" rectangle [anchor,1]ᵈ.", {})]),
    (0, [("fig 1", {"bold": True, "color": ACCENT}),
         (" — a single rectangular corner: the rectangle = the whole positive region. ",
          {}), ("Exact.", {"bold": True, "color": GREEN})]),
    (0, [("fig 2", {"bold": True, "color": ACC2}),
         (" — a ", {}), ("staircase", {"bold": True}),
         (" frontier (≥2 corners; happens when the predicates are correlated / the joint isn't axis-separable).", {})]),
    (0, [("One rectangle can't hug every step → positive points in the ", {}),
         ("other steps", {"bold": True}),
         (" (low-s₁/high-s₂ and high-s₁/low-s₂) fall off the grid and are ", {}),
         ("never sampled.", {"bold": True, "color": ACC2})]),
    (0, [("Those near-frontier points are exactly where joint cardinality just ", {}),
         ("switches on", {"italic": True}), (" → ", {}),
         ("plan flips & q-error cliffs", {"bold": True}), (" — the interesting ones.", {})]),
])
chip(s, Inches(0.4), Inches(6.45), Inches(7.55),
     [("no dead points are added — the rectangle stays inside the positive region — but some interesting near-frontier points are lost.",
       {"bold": True})], ACC2, size=11, h=0.55)
textbox(s, Inches(8.15), Inches(6.45), Inches(4.65), Inches(0.6),
        [[("Mitigation: detect a multi-corner frontier and union per-step rectangles (a staircase-aware lower envelope) instead of one anchor.",
           {"size": 10.5, "italic": True, "color": GREY})]])

# 8 — summary
s = slide("Summary", "One anchor, computed once per query, fixes the boundary")
bullets(s, LX, Inches(1.6), Inches(12.23), Inches(4.6), [
    (0, [("Problem:", {"bold": True, "color": RED}), (" a grid over ", {}),
         ("[global min, max]", {"font": FM}), (" wastes its low-selectivity corner on 0-row joints.", {})]),
    (0, [("Fix:", {"bold": True, "color": GREEN}), ("  ", {}),
         ("find_min_positive_point_nd", {"font": FM, "bold": True}),
         (" finds the per-axis minimum selectivities whose joint cardinality is > 0.", {})]),
    (1, [("diagonal doubling", {"bold": True}), (" → jump onto the positive region (geometric, cheap),", {})]),
    (1, [("binary search", {"bold": True}), (" → pin the diagonal crossing s*,", {})]),
    (1, [("coordinate descent", {"bold": True}), (" → shrink the square to the tight per-axis anchor.", {})]),
    (0, [("Effect:", {"bold": True}), (" the grid is re-based to ", {}),
         ("[anchor, max]", {"font": FM, "bold": True}),
         (" — every one of the 11×11 points returns rows; no dead boundary, no wasted profiling.", {})]),
    (0, [("Cost:", {"bold": True}), (" a handful of COUNT(*) probes per query, run once before sampling.", {})]),
])

# =========================================================
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "grid_boundary_anchoring.pptx")
prs.save(OUT)
print("saved", os.path.abspath(OUT), "slides:", _page[0], "| dead(old grid):", N_DEAD)
