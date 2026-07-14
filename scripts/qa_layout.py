#!/usr/bin/env python3
"""Static layout QA for the generated pptx (no renderer available)."""
import sys
from pptx import Presentation
from pptx.util import Emu
EMU=914400.0
SW,SH=13.333,7.5
FOOT=7.05
def inch(v): return v/EMU
p=Presentation(sys.argv[1] if len(sys.argv)>1 else "../interpolation_methods_explained.pptx")
issues=0
for si,sl in enumerate(p.slides,1):
    boxes=[]
    for sh in sl.shapes:
        try:
            l,t,w,h=inch(sh.left),inch(sh.top),inch(sh.width),inch(sh.height)
        except Exception:
            continue
        r,b=l+w,t+h
        tag=sh.shape_type
        # off-slide
        if l< -0.02 or t< -0.02 or r>SW+0.02 or b>SH+0.02:
            print(f"[S{si:2d}] OFF-SLIDE  {str(tag):22s} L{l:.2f} T{t:.2f} R{r:.2f} B{b:.2f}")
            issues+=1
        # table/box into footer
        if sh.has_table and b>FOOT:
            print(f"[S{si:2d}] TABLE→FOOTER  bottom={b:.2f} (>{FOOT})")
            issues+=1
        # long mono lines (Consolas ~0.106in per char at given pt; check 14pt boxes)
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                txt="".join(r.text for r in para.runs)
                sz=max([(r.font.size.pt if r.font.size else 14) for r in para.runs]+[14])
                # only check monospace-ish (Consolas) runs
                mono_=bool(txt) and all((r.font.name=="Consolas") for r in para.runs if r.text.strip())
                if mono_ and txt:
                    cw=0.0073*sz  # approx Consolas char width in inches per pt-size
                    need=len(txt)*cw
                    if need>w+0.08:
                        print(f"[S{si:2d}] MONO OVERFLOW ~{need:.2f}in > box {w:.2f}in  | {txt[:46]!r} sz{sz}")
                        issues+=1

    # --- block overlap check (tables vs rounded-rect panels/chips) ---
    blocks=[]
    for sh in sl.shapes:
        try:
            l,t,w,h=inch(sh.left),inch(sh.top),inch(sh.width),inch(sh.height)
        except Exception: continue
        is_tbl=sh.has_table
        is_panel=False
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            if sh.shape_type==MSO_SHAPE_TYPE.AUTO_SHAPE and w*h>0.25 and h>0.3:
                is_panel=True
        except Exception: pass
        if is_tbl or is_panel:
            blocks.append((("T" if is_tbl else "P"),l,t,l+w,t+h))
    for a in range(len(blocks)):
        for b in range(a+1,len(blocks)):
            ka,la,ta,ra,ba=blocks[a]; kb,lb,tb,rb,bb=blocks[b]
            ix=max(0,min(ra,rb)-max(la,lb)); iy=max(0,min(ba,bb)-max(ta,tb))
            area=ix*iy
            # ignore panel-contains (a textbox sits ON its own panel is fine, but those are textboxes not blocks)
            if area>0.15 and not (ka=="P" and kb=="P" and area>0.9*min((ra-la)*(ba-ta),(rb-lb)*(bb-tb))):
                print(f"[S{si:2d}] OVERLAP {ka}+{kb} area={area:.2f}in^2  @({max(la,lb):.1f},{max(ta,tb):.1f})")
                issues+=1


    # --- table/chip vs textbox collisions ---
    tbls=[]; chips=[]; tboxes=[]
    from pptx.enum.shapes import MSO_SHAPE_TYPE as MST
    for sh in sl.shapes:
        try: l,t,w,h=inch(sh.left),inch(sh.top),inch(sh.width),inch(sh.height)
        except Exception: continue
        box=(l,t,l+w,t+h)
        if sh.has_table: tbls.append(box)
        elif sh.shape_type==MST.AUTO_SHAPE and h<0.7 and w>1.0:
            # chip-like rounded rect (has its own centered text)
            if sh.has_text_frame and sh.text_frame.text.strip(): chips.append(box)
        elif sh.has_text_frame and sh.text_frame.text.strip():
            tboxes.append(box)
    def ov(a,b):
        ix=max(0,min(a[2],b[2])-max(a[0],b[0])); iy=max(0,min(a[3],b[3])-max(a[1],b[1]))
        return ix*iy
    for tb in tbls:
        for x in tboxes:
            if ov(tb,x)>0.08:
                print(f"[S{si:2d}] TABLE×TEXT overlap {ov(tb,x):.2f}in^2 @({max(tb[0],x[0]):.1f},{max(tb[1],x[1]):.1f})")
                issues+=1
    for c in chips:
        for x in tboxes:
            if ov(c,x)>0.10:
                print(f"[S{si:2d}] CHIP×TEXT overlap {ov(c,x):.2f}in^2 @({max(c[0],x[0]):.1f},{max(c[1],x[1]):.1f})")
                issues+=1

print(f"\n{'OK - no issues' if issues==0 else str(issues)+' issue(s)'}  across {len(p.slides._sldIdLst)} slides")
